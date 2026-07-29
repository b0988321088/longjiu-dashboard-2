"""
策略：只取代為長度「小於等於」原文的文字，避免溢出
超過長度的文字不取代，保持原樣
"""
from pptx import Presentation

prs = Presentation('龍九簡報模板.pptx')

repl = {
    '轉貸投資可行性評估計劃': '資產日跌20萬對策',
    '以低利資金 2.185% 優化負債結構 × 提升被動現金流 +106K/月': '2026-07-28 台美股重挫 證券-98K 基金-44K',
    '2026 年 7 月 26 日 ｜ 評估期間：2026 Q3–Q4': '2026-07-28 ｜ Chief Secretary 緊急應變',
    'CONTENTS': '應變對策',
    '01 · 核心問題：每月現金流缺口 -4,099': '01 · 今日跌 -3.77%/-5.58% = -14.2萬',
    '02 · 資產負債結構：負債率 41.8%': '02 · 穿透：台股9.4%遠低目標35%',
    '03 · 轉貸方案：築巢優利貸 2.185% 為最優解': '03 · 三種策略：不動/減碼美股/加碼台股',
    '04 · 三階段部署：清償 → 安全網 → 擴張': '04 · 三階段：觀察→決策→執行',
    '05 · 第一階段：清償 4M 高利負債，月省 16K': '05 · 減碼：美股34%降5pp',
    '06 · 第二階段：保留 3M 額度，建立緩衝': '06 · 加碼：台股9.4%補5pp',
    '07 · 第三階段：低利套利，目標報酬 >6%': '07 · 防禦：現金311萬+00983D',
    '08 · ETF 建倉策略：00983D/00919/00878': '08 · 三軌並行：防禦50+加碼80+備援',
    '09 · 保單第三站：PIMCO+AI+A10 月配 17.5K': '09 · 明日開盤行動：跌>2%進20萬',
    '10 · 執行時間表：7 月至 12 月': '10 · 壓力測試：再跌5/10/20%',
    '11 · 預期成效：月淨現金流 +102,259': '11 · 風險矩陣：5項風險全可控',
    '12 · 風險分析與緩解措施': '12 · 結論：不恐慌按紀律',
    '13 · 結論：每月從 -4,099 到 +102,259': '13 · 總結：系統因應',
}

count = 0
for slide in prs.slides:
    for shape in slide.shapes:
        if shape.has_text_frame:
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    o = run.text
                    n = o
                    for old_t, new_t in repl.items():
                        if old_t in n and len(new_t) <= len(old_t):
                            n = n.replace(old_t, new_t)
                    if n != o and len(n) <= len(o):
                        run.text = n
                        count += 1
                    elif n != o:
                        # Try to shorten replacement
                        short_t = new_t[:len(old_t)]
                        if short_t != o:
                            run.text = short_t
                            count += 1

prs.save('資產應變對策.pptx')
print(f'Done - {count} updates (length-safe)')
