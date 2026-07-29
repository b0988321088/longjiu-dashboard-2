from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import MSO_ANCHOR, MSO_AUTO_SIZE
from pptx.enum.shapes import MSO_SHAPE # Import MSO_SHAPE
from pptx.dml.color import RGBColor # Import RGBColor
import os

def create_adaptive_presentation(title="龍九控股自適應簡報", date="2026年7月29日"):
    prs = Presentation()

    # --- Slide 1: Cover Page ---
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    title_placeholder = slide.shapes.title
    subtitle_placeholder = slide.placeholders[1]

    title_placeholder.text = title
    subtitle_placeholder.text = f"報告日期: {date}\n自適應內容引擎演示"

    # --- Slide 2: 市場概況與核心觀察 (Content Slide) ---
    blank_slide_layout = prs.slide_layouts[5]
    slide = prs.slides.add_slide(blank_slide_layout)

    # Title
    left, top, width, height = Inches(1), Inches(0.5), Inches(8), Inches(1)
    title_shape = slide.shapes.add_textbox(left, top, width, height)
    text_frame = title_shape.text_frame
    p = text_frame.paragraphs[0]
    p.text = "市場概況與核心觀察"
    p.font.size = Pt(36)
    p.font.bold = True

    # Content Placeholder (adaptive)
    left, top, width, height = Inches(1), Inches(1.5), Inches(8), Inches(5.5)
    content_shape = slide.shapes.add_textbox(left, top, width, height)
    text_frame = content_shape.text_frame
    text_frame.word_wrap = True
    text_frame.auto_size = MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT
    
    p = text_frame.paragraphs[0]
    p.text = "美國10年期公債殖利率：目前為4.602%，接近52週高點4.71%，遠高於2020年低點0.32%，顯示債券市場發生結構性反轉。資金成本顯著上升。"
    p.font.size = Pt(18)
    
    p = text_frame.add_paragraph()
    p.text = "股市表現：道瓊工業指數上漲0.59%至52,517點，主要由防禦性板塊（如工業、金融、公用事業）領漲。標準普爾500指數微跌0.14%至7,403點，納斯達克指數下跌0.68%至24,763點，顯示科技股近期表現疲軟，市場資金正從高成長科技股輪動至價值型和防禦性板塊。"
    p.font.size = Pt(18)

    p = text_frame.add_paragraph()
    p.text = "外匯市場：美元兌新台幣匯率為29.78，台幣呈現偏弱走勢，可能受到美國升息預期及資金外流壓力影響。"
    p.font.size = Pt(18)

    # --- Slide 3: 要點列表 (Bullet Points Slide) ---
    slide = prs.slides.add_slide(blank_slide_layout)

    # Title
    left, top, width, height = Inches(1), Inches(0.5), Inches(8), Inches(1)
    title_shape = slide.shapes.add_textbox(left, top, width, height)
    text_frame = title_shape.text_frame
    p = text_frame.paragraphs[0]
    p.text = "本週重點摘要"
    p.font.size = Pt(36)
    p.font.bold = True

    # Bullet Points Placeholder (adaptive)
    left, top, width, height = Inches(1), Inches(1.5), Inches(8), Inches(6)
    bullet_shape = slide.shapes.add_textbox(left, top, width, height)
    text_frame = bullet_shape.text_frame
    text_frame.word_wrap = True
    text_frame.auto_size = MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT

    p = text_frame.paragraphs[0]
    p.text = "全球經濟成長放緩預期，投資者轉向避險資產。"
    p.level = 0
    p.font.size = Pt(20)

    p = text_frame.add_paragraph()
    p.text = "主要央行貨幣政策分歧，影響匯率波動。"
    p.level = 0
    p.font.size = Pt(20)

    p = text_frame.add_paragraph()
    p.text = "新興市場面臨資本外流壓力，需謹慎評估投資風險。"
    p.level = 0
    p.font.size = Pt(20)

    p = text_frame.add_paragraph()
    p.text = "能源價格持續高位運行，通脹壓力依然存在。"
    p.level = 0
    p.font.size = Pt(20)

    # --- Slide 4: 數據卡片比較 (Data Card Comparison Slide) ---
    slide = prs.slides.add_slide(blank_slide_layout)

    # Title
    left, top, width, height = Inches(1), Inches(0.5), Inches(8), Inches(1)
    title_shape = slide.shapes.add_textbox(left, top, width, height)
    text_frame = title_shape.text_frame
    p = text_frame.paragraphs[0]
    p.text = "關鍵市場數據一覽"
    p.font.size = Pt(36)
    p.font.bold = True

    # Market Data
    market_data = [
        {"title": "US 10yr Treasury", "value": "4.602%", "detail": "52週高:4.71%, 低:3.93%"},
        {"title": "道瓊指數", "value": "52,517", "detail": "+0.59% (防禦股領漲)"},
        {"title": "S&P 500", "value": "7,403", "detail": "-0.14% (盤整)"},
        {"title": "Nasdaq", "value": "24,763", "detail": "-0.68% (科技弱)"},
        {"title": "USD/TWD", "value": "29.78", "detail": "台幣偏弱"}
    ]

    total_cards = len(market_data)
    available_width = prs.slide_width - Inches(1) # Total slide width minus left/right margins
    spacing = Inches(0.2)
    card_width_value = (available_width - (spacing * (total_cards - 1))) / total_cards

    card_height = Inches(3)
    start_left = Inches(0.5)
    card_top = Inches(2)

    for i, data in enumerate(market_data):
        left = start_left + (card_width_value + spacing) * i
        
        # Add a rectangle as a "card" background
        shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, card_top, card_width_value, card_height)
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(0x00, 0x7B, 0xFF) # Blue color

        # Add text box for card content
        text_box = slide.shapes.add_textbox(left + Inches(0.1), card_top + Inches(0.1), card_width_value - Inches(0.2), card_height - Inches(0.2))
        text_frame = text_box.text_frame
        text_frame.word_wrap = True
        text_frame.auto_size = MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT
        text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE # Center text vertically

        # Title
        p = text_frame.paragraphs[0]
        p.text = data["title"]
        p.font.size = Pt(20)
        p.font.bold = True
        p.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF) # White text

        # Value
        p = text_frame.add_paragraph()
        p.text = data["value"]
        p.font.size = Pt(30)
        p.font.bold = True
        p.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)

        # Detail
        p = text_frame.add_paragraph()
        p.text = data["detail"]
        p.font.size = Pt(14)
        p.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
        
    # --- Slide 5: Conclusion Page ---
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    title_placeholder = slide.shapes.title
    subtitle_placeholder = slide.placeholders[1]

    title_placeholder.text = "結論與展望"
    subtitle_placeholder.text = "市場輪動趨勢明確，防禦性資產配置重要性提升。"

    prs.save(os.path.join(os.path.dirname(os.path.abspath(__file__)), "adaptive_presentation_demo.pptx"))
    print(f"簡報 'adaptive_presentation_demo.pptx' 已更新成功。")

if __name__ == "__main__":
    create_adaptive_presentation()
