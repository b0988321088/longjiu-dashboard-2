"""
slide_engine.py — 龍九簡報引擎 v3
用法：
  python slide_engine.py content.json
"""
import json, sys
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

C = {
    'bg': RGBColor(0x0B,0x0D,0x1A),
    'card': RGBColor(0x15,0x18,0x2E),
    'gold': RGBColor(0xF7,0xA0,0x1C),
    'green': RGBColor(0x34,0xD3,0x99),
    'red': RGBColor(0xFF,0x5C,0x5C),
    'blue': RGBColor(0x5B,0x9B,0xF7),
    'text': RGBColor(0xE8,0xE8,0xF0),
    'gray': RGBColor(0x8A,0x8F,0xA0),
    'white': RGBColor(0xFF,0xFF,0xFF),
}

def new_ppt():
    p = Presentation()
    p.slide_width = Inches(13.33); p.slide_height = Inches(7.5)
    return p

def new_slide(p):
    s = p.slides.add_slide(p.slide_layouts[6])
    s.background.fill.solid(); s.background.fill.fore_color.rgb = C['bg']
    return s

def box(s, l, t, w, h):
    sh = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(l), Inches(t), Inches(w), Inches(h))
    sh.fill.solid(); sh.fill.fore_color.rgb = C['card']
    sh.line.color.rgb = RGBColor(0x1E,0x22,0x3D); sh.line.width = Pt(1)
    return sh

def add_text(s, text, l, t, w, h, sz=14, bold=False, clr='text', align='l'):
    c = C.get(clr, C['text'])
    a = {'l':PP_ALIGN.LEFT,'c':PP_ALIGN.CENTER,'r':PP_ALIGN.RIGHT}.get(align, PP_ALIGN.LEFT)
    tb = s.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = True
    tf.auto_size = None  # 不自動調整框大小
    # 設定垂直置中 + 允許縮小
    from pptx.oxml.ns import qn
    bodyPr = tf._txBody.bodyPr
    bodyPr.set('autofit', 'shrink')
    p = tf.paragraphs[0]; p.text = text
    p.font.size = Pt(sz); p.font.bold = bold; p.font.color.rgb = c; p.alignment = a

def add_lines(s, items, l, t, w, h, sz=13, clr='text'):
    c = C.get(clr, C['text'])
    tb = s.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tb.text_frame.word_wrap = True
    for i, item in enumerate(items):
        p = tb.text_frame.paragraphs[0] if i==0 else tb.text_frame.add_paragraph()
        p.text = item; p.font.size = Pt(sz); p.font.color.rgb = c; p.space_after = Pt(3)

def divider(s, l, t, w):
    box(s, l, t, w, 0.04)

def metric(s, l, t, w, h, label, val, sub='', clr='green'):
    box(s, l, t, w, h)
    add_text(s, label, l+0.2, t+0.15, w-0.4, 0.3, 11, clr='gray')
    add_text(s, val, l+0.2, t+0.5, w-0.4, 0.5, 26, True, clr)
    if sub: add_text(s, sub, l+0.2, t+1.15, w-0.4, 0.3, 10, clr='gray')

# ====== Slide builders ======

def build_cover(s, cfg):
    box(s, 1, 1.5, 11.3, 4.5)
    add_text(s, cfg.get('brand','龍九控股'), 1.5, 1.8, 5, 0.5, 14, clr='blue')
    add_text(s, cfg['title'], 1.5, 2.3, 10, 1.2, 40, True, 'white')
    add_text(s, cfg.get('subtitle',''), 1.5, 3.5, 10, 0.5, 16, clr='gray')
    add_text(s, cfg.get('date',''), 1.5, 4.5, 8, 0.4, 12, clr='gray')

def build_toc(s, cfg):
    add_text(s, '目錄', 1, 0.5, 6, 0.8, 32, True, 'white')
    divider(s, 1, 1.3, 3)
    for i, item in enumerate(cfg.get('items',[])):
        add_text(s, item, 1.5, 2+i*0.38, 10, 0.35, 14)

def build_metrics(s, cfg):
    add_text(s, cfg['title'], 0.8, 0.4, 11, 0.6, 26, True, 'white')
    if 'subtitle' in cfg: add_text(s, cfg['subtitle'], 0.8, 1.1, 11, 0.4, 14, clr='gray')
    divider(s, 0.8, 1.2, 3)
    for i, m in enumerate(cfg.get('metrics',[])):
        metric(s, 0.5+i*2.4, 1.6, 2.2, 1.6, m.get('label',''), m.get('val',''), m.get('sub',''), m.get('color','green'))
    lc = cfg.get('left_card')
    if lc:
        box(s, 0.5, 3.5, 6, 2.5)
        add_text(s, lc.get('title',''), 0.8, 3.7, 5.4, 0.4, 16, True, lc.get('tc','white'))
        add_lines(s, lc.get('items',[]), 0.8, 4.2, 5.4, 1.5, 13)
    rc = cfg.get('right_card')
    if rc:
        box(s, 6.8, 3.5, 6, 2.5)
        add_text(s, rc.get('title',''), 7.1, 3.7, 5.4, 0.4, 16, True, rc.get('tc','white'))
        add_lines(s, rc.get('items',[]), 7.1, 4.2, 5.4, 1.5, 13)

def build_compare(s, cfg):
    add_text(s, cfg['title'], 0.8, 0.4, 11, 0.6, 26, True, 'white')
    add_text(s, cfg.get('subtitle',''), 0.8, 1.1, 11, 0.4, 14, clr='gray')
    for i, p in enumerate(cfg.get('plans',[])):
        l = 0.5 + i*4.2
        box(s, l, 1.6, 3.9, 3.5)
        add_text(s, p.get('name',''), l+0.2, 1.8, 3.5, 0.3, 16, True, p.get('color','white'))
        add_text(s, p.get('rate',''), l+0.2, 2.2, 3.5, 0.3, 24, True, 'white')
        add_lines(s, p.get('items',[]), l+0.2, 2.7, 3.5, 2, 11)
    cc = cfg.get('conclusion')
    if cc:
        box(s, 0.5, 5.4, 12.3, 1.5)
        add_text(s, cc.get('title','結論'), 0.8, 5.6, 11.7, 0.4, 16, True, 'gold')
        add_lines(s, cc.get('items',[]), 0.8, 6.1, 11.7, 0.8, 13)

def build_phases(s, cfg):
    add_text(s, cfg['title'], 0.8, 0.4, 11, 0.6, 26, True, 'white')
    for i, p in enumerate(cfg.get('phases',[])):
        l = 0.5 + i*4.2
        box(s, l, 1.5, 3.9, 4)
        add_text(s, p.get('period',''), l+0.2, 1.7, 3.5, 0.3, 12, clr=p.get('color','gray'))
        add_text(s, p.get('title',''), l+0.2, 2.1, 3.5, 0.3, 16, True, 'white')
        add_lines(s, p.get('items',[]), l+0.2, 2.6, 3.5, 2.5, 12)
    if 'conclusion' in cfg:
        box(s, 0.5, 5.8, 12.3, 1.2)
        add_text(s, cfg['conclusion'], 0.8, 6.0, 11.7, 0.8, 15, True, 'green')

def build_risks(s, cfg):
    add_text(s, cfg['title'], 0.8, 0.4, 11, 0.6, 26, True, 'white')
    for i, r in enumerate(cfg.get('items',[])):
        l = 0.5 + i*2.5
        box(s, l, 1.3, 2.3, 3)
        add_text(s, r.get('name',''), l+0.2, 1.5, 1.9, 0.3, 13, True, r.get('color','text'))
        add_lines(s, r.get('desc',[]), l+0.2, 1.9, 1.9, 1.5, 11)

BUILDERS = {
    'cover': build_cover, 'toc': build_toc,
    'metrics': build_metrics, 'compare': build_compare,
    'phases': build_phases, 'risks': build_risks,
}

def build(config_file):
    cfg = json.load(open(config_file, encoding='utf-8'))
    p = new_ppt()
    for sc in cfg['slides']:
        s = new_slide(p)
        fn = BUILDERS.get(sc['type'])
        if fn: fn(s, sc)
    out = cfg.get('output', 'output.pptx')
    p.save(out)
    print(f'OK {out} ({len(cfg["slides"])} slides)')

if __name__ == '__main__':
    build(sys.argv[1] if len(sys.argv)>1 else 'content.json')
