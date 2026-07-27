"""
龍九控股 簡報引擎 (Presentation Engine)
==========================================
基於 professional-presentation 技能框架，支援所有未來簡報需求。
遵循 SCQA 故事線、行動標題、12 欄網格、卡片化設計系統。
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.enum.chart import XL_CHART_TYPE
from pptx.dml.color import RGBColor
from pptx.chart.data import CategoryChartData

# ═══════════════════════════════════════
# 設計系統 (Design System) — 不可修改
# ═══════════════════════════════════════

# 主色（60%）
BG = RGBColor(0x0B, 0x0D, 0x1A)          # 深色背景
CARD = RGBColor(0x15, 0x18, 0x2E)         # 卡片底色
CARD_BORDER = RGBColor(0x1E, 0x22, 0x3D)  # 卡片邊框
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

# 輔助色（30%）
TEXT = RGBColor(0xE8, 0xE8, 0xF0)         # 主文字
GRAY = RGBColor(0x8A, 0x8F, 0xA0)         # 輔助文字

# 強調色（10%）— 單頁不超過 2 個
GOLD = RGBColor(0xF7, 0xA0, 0x1C)         # 重點標示
GREEN = RGBColor(0x34, 0xD3, 0x99)        # 正向/收益
RED = RGBColor(0xFF, 0x5C, 0x5C)          # 負向/警示
BLUE = RGBColor(0x5B, 0x9B, 0xF7)         # 資訊
PURPLE = RGBColor(0xA7, 0x8B, 0xFA)       # 次要區分

# 網格系統（12 欄）
COL = Inches(13.33) / 12  # ~1.11 吋/欄
ROW_H = Inches(0.7)       # 標準卡片高度


class SlideDeck:
    """簡報引擎 — 所有未來簡報皆使用此類別"""

    def __init__(self, title="龍九控股簡報", subtitle=""):
        self.prs = Presentation()
        self.prs.slide_width = Inches(13.33)
        self.prs.slide_height = Inches(7.5)

    # ── 基礎元件 ──

    def _new(self):
        """建立新投影片"""
        sl = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        sl.background.fill.solid()
        sl.background.fill.fore_color.rgb = BG
        return sl

    def card(self, sl, left, top, w, h):
        """加入卡片（圓角）"""
        s = sl.shapes.add_shape(1, Inches(left), Inches(top), Inches(w), Inches(h))
        s.fill.solid()
        s.fill.fore_color.rgb = CARD
        s.line.color.rgb = CARD_BORDER
        s.line.width = Pt(1)
        return s

    def txt(self, sl, text, left, top, w, h, size=14, bold=False, color=TEXT, align=PP_ALIGN.LEFT):
        """加入文字方塊"""
        tb = sl.shapes.add_textbox(Inches(left), Inches(top), Inches(w), Inches(h))
        tb.text_frame.word_wrap = True
        p = tb.text_frame.paragraphs[0]
        p.text = text
        p.font.size = Pt(size)
        p.font.bold = bold
        p.font.color.rgb = color
        p.alignment = align
        return tb

    def metric_card(self, sl, left, top, w, h, label, value, highlight=GREEN, sub=""):
        """加入指標大數字卡片（Metric Card）"""
        self.card(sl, left, top, w, h)
        self.txt(sl, label, left + 0.2, top + 0.15, w - 0.4, 0.3, 12, False, GRAY)
        self.txt(sl, str(value), left + 0.2, top + 0.5, w - 0.4, 0.6, 36, True, highlight)
        if sub:
            self.txt(sl, sub, left + 0.2, top + h - 0.4, w - 0.4, 0.3, 11, False, GRAY)

    def page_title(self, sl, title, action_line=""):
        """加入頁面標題（Action Title 格式）"""
        self.txt(sl, title, 0.8, 0.3, 11, 0.7, 30, True, WHITE)
        # 金線
        self.card(sl, 0.8, 1.0, 3, 0.03)
        if action_line:
            self.txt(sl, action_line, 0.8, 1.2, 11, 0.4, 14, False, GRAY)

    def multi(self, sl, lines, left, top, w, h, size=13, color=TEXT):
        """多行文字"""
        tb = sl.shapes.add_textbox(Inches(left), Inches(top), Inches(w), Inches(h))
        tb.text_frame.word_wrap = True
        for i, line in enumerate(lines):
            p = tb.text_frame.paragraphs[0] if i == 0 else tb.text_frame.add_paragraph()
            p.text = line
            p.font.size = Pt(size)
            p.font.color.rgb = color
            p.space_after = Pt(4)

    # ── 投影片類型 ──

    def slide_cover(self, title, subtitle="", date="", author=""):
        """封面"""
        sl = self._new()
        self.txt(sl, "龍九控股", 1, 1, 5, 0.6, 16, False, GRAY)
        self.txt(sl, title, 1, 1.8, 11, 1.5, 44, True, WHITE)
        self.card(sl, 1, 3.5, 4, 0.04)
        if subtitle:
            self.txt(sl, subtitle, 1, 3.9, 11, 0.6, 20, False, GRAY)
        if date:
            self.txt(sl, date, 1, 5.5, 8, 0.4, 14, False, GRAY)
        return sl

    def slide_toc(self, title, items):
        """目錄"""
        sl = self._new()
        self.txt(sl, title, 0.8, 0.4, 8, 0.6, 28, True, WHITE)
        self.card(sl, 0.8, 1.1, 3, 0.03)
        for i, item in enumerate(items):
            self.txt(sl, f"{'%02d' % (i+1)} · {item}", 1.5, 1.8 + i * 0.4, 10, 0.35, 15)
        return sl

    def slide_metrics(self, title, action_line, metrics, cols=3):
        """指標總覽頁（Metric Cards 陣列）"""
        sl = self._new()
        self.page_title(sl, title, action_line)
        card_w = (11.5 / cols)
        for i, (label, value, highlight, sub) in enumerate(metrics):
            col = i % cols
            row = i // cols
            x = 0.8 + col * (card_w + 0.3)
            y = 1.8 + row * 1.6
            self.metric_card(sl, x, y, card_w, 1.3, label, value, highlight, sub)
        return sl

    def slide_comparison(self, title, action_line, left_data, right_data):
        """對比頁（雙欄）"""
        sl = self._new()
        self.page_title(sl, title, action_line)
        cols = [left_data, right_data]
        for i, (header, items, color) in enumerate(cols):
            x = 0.5 + i * 6.3
            self.card(sl, x, 1.8, 5.8, 5)
            self.txt(sl, header, x + 0.3, 2, 5, 0.5, 22, True, color)
            for j, (label, value, vcolor) in enumerate(items):
                y = 2.8 + j * 0.7
                self.txt(sl, label, x + 0.3, y, 3, 0.4, 13, False, GRAY)
                self.txt(sl, value, x + 3.5, y, 2, 0.4, 15, True, vcolor)
        return sl

    def slide_timeline(self, title, action_line, events):
        """時間軸頁"""
        sl = self._new()
        self.page_title(sl, title, action_line)
        for i, (phase, event, color) in enumerate(events):
            y = 1.8 + i * 0.72
            self.card(sl, 0.5, y, 12, 0.6)
            self.txt(sl, phase, 0.8, y + 0.08, 2.5, 0.4, 14, True, color)
            self.txt(sl, event, 3.5, y + 0.08, 8, 0.4, 14)
        return sl

    def slide_risks(self, title, action_line, risks):
        """風險分析頁"""
        sl = self._new()
        self.page_title(sl, title, action_line)
        for i, (risk, impact, mitigation) in enumerate(risks):
            y = 1.8 + i * 1
            self.card(sl, 0.5, y, 12, 0.8)
            self.txt(sl, risk, 0.8, y + 0.1, 3, 0.4, 16, True, RED if '風險' in risk else GOLD)
            self.txt(sl, impact, 4.2, y + 0.1, 4, 0.4, 12, False, GRAY)
            self.txt(sl, mitigation, 8.5, y + 0.1, 4, 0.4, 12, False, GREEN)
        return sl

    def slide_summary(self, title, items):
        """總結頁"""
        sl = self._new()
        self.txt(sl, title, 0.8, 0.4, 8, 0.6, 28, True, WHITE)
        self.card(sl, 0.8, 1.1, 3, 0.03)
        for i, (icon, heading, desc, color) in enumerate(items):
            y = 1.5 + i * 1
            self.card(sl, 0.5, y, 12, 0.8)
            self.txt(sl, f"{icon}  {heading}", 0.8, y + 0.1, 3, 0.4, 20, True, color)
            self.txt(sl, desc, 4.2, y + 0.15, 8, 0.4, 14)
        return sl

    def save(self, path):
        self.prs.save(path)
        return path

    # ── 新增版型：多樣化布局 ──

    def slide_section(self, number, title, subtitle=""):
        """章節分隔頁（大字＋編號）"""
        sl = self._new()
        self.txt(sl, f"{number:02d}", 1, 1.5, 3, 2, 72, True, GOLD)
        self.card(sl, 1, 3.8, 4, 0.04)
        self.txt(sl, title, 1, 4.2, 11, 1, 36, True, WHITE)
        if subtitle:
            self.txt(sl, subtitle, 1, 5.5, 10, 0.6, 18, False, GRAY)
        return sl

    def slide_2col(self, title, action_line, left_title, left_items, right_title, right_items):
        """雙欄對比（不同底色區隔）"""
        sl = self._new()
        self.page_title(sl, title, action_line)
        # 左欄
        self.card(sl, 0.5, 1.8, 5.8, 5)
        self.txt(sl, left_title, 0.8, 2, 5, 0.5, 22, True, GOLD)
        for i, (label, value, color) in enumerate(left_items):
            y = 2.7 + i * 0.65
            self.txt(sl, label, 0.8, y, 3, 0.4, 13, False, GRAY)
            self.txt(sl, value, 4, y, 2, 0.4, 15, True, color)
        # 右欄
        bg2 = RGBColor(0x0F, 0x12, 0x1E)
        s = sl.shapes.add_shape(1, Inches(6.8), Inches(1.8), Inches(5.8), Inches(5))
        s.fill.solid(); s.fill.fore_color.rgb = bg2
        s.line.color.rgb = CARD_BORDER; s.line.width = Pt(1)
        self.txt(sl, right_title, 7.1, 2, 5, 0.5, 22, True, GREEN)
        for i, (label, value, color) in enumerate(right_items):
            y = 2.7 + i * 0.65
            self.txt(sl, label, 7.1, y, 3, 0.4, 13, False, GRAY)
            self.txt(sl, value, 10.5, y, 2, 0.4, 15, True, color)
        return sl

    def slide_process(self, title, action_line, steps):
        """流程圖頁（橫向箭頭）"""
        sl = self._new()
        self.page_title(sl, title, action_line)
        n = len(steps)
        w = min(3.5, 11 / n)
        for i, (label, desc, color) in enumerate(steps):
            x = 0.5 + i * (w + 0.3)
            self.card(sl, x, 2, w, 3.5)
            self.card(sl, x, 2, w, 0.06)
            # 大數字
            self.txt(sl, f"Step {i+1}", x + 0.2, 2.3, w - 0.4, 0.4, 14, False, color)
            self.txt(sl, label, x + 0.2, 2.8, w - 0.4, 0.6, 24, True, color)
            self.txt(sl, desc, x + 0.2, 3.7, w - 0.4, 1.5, 13, False, GRAY)
            # 箭頭
            if i < n - 1:
                self.txt(sl, "→", x + w - 0.1, 3, 0.5, 0.5, 24, True, GRAY)
        return sl

    def slide_quote(self, quote, author, title=""):
        """引言/重點提示頁"""
        sl = self._new()
        self.txt(sl, '"', 1, 1.5, 2, 1.5, 96, True, GOLD)
        self.txt(sl, quote, 2, 2.5, 10, 3, 28, True, WHITE)
        self.txt(sl, f"— {author}", 7, 5.5, 5, 0.5, 16, False, GRAY)
        if title:
            self.txt(sl, title, 7, 6, 5, 0.4, 13, False, GOLD)
        return sl

    def slide_numbers(self, title, action_line, big_numbers):
        """數據故事頁（大數字為主角）"""
        sl = self._new()
        self.page_title(sl, title, action_line)
        for i, (number, label, sub, color) in enumerate(big_numbers):
            x = 0.5 + i * 4.2
            self.card(sl, x, 2.5, 3.8, 3)
            self.txt(sl, str(number), x + 0.3, 2.8, 3.2, 1, 56, True, color)
            self.txt(sl, label, x + 0.3, 3.8, 3.2, 0.4, 18, True, WHITE)
            if sub:
                self.txt(sl, sub, x + 0.3, 4.3, 3.2, 0.4, 13, False, GRAY)
        return sl


# ═══════════════════════════════════════
# 轉貸投資簡報（使用簡報引擎）
# ═══════════════════════════════════════

if __name__ == "__main__":
    d = SlideDeck()

    # S1: 封面
    d.slide_cover(
        title="轉貸投資可行性評估計劃",
        subtitle="以低利資金 2.185% 優化負債結構 × 提升被動現金流 +106K/月",
        date="2026 年 7 月 26 日 ｜ 評估期間：2026 Q3–Q4",
    )

    # S2: 目錄
    d.slide_toc("CONTENTS", [
        "核心問題：每月現金流缺口 -4,099",
        "資產負債結構：負債率 41.8%",
        "轉貸方案：築巢優利貸 2.185% 為最優解",
        "三階段部署：清償 → 安全網 → 擴張",
        "第一階段：清償 4M 高利負債，月省 16K",
        "第二階段：保留 3M 額度，建立緩衝",
        "第三階段：低利套利，目標報酬 >6%",
        "ETF 建倉策略：00983D/00919/00878",
        "保單第三站：PIMCO+AI+A10 月配 17.5K",
        "執行時間表：7 月至 12 月",
        "預期成效：月淨現金流 +102,259",
        "風險分析與緩解措施",
        "結論：每月從 -4,099 到 +102,259",
    ])

    # S3: 核心問題（SCQA：Complication）
    sl = d._new()
    d.page_title(sl, "每月現金流缺口 -4,099，保單利息吃掉 16,000", "保單借貸利率 ~5% 為最高成本負債，每月利息相當於一筆房貸月付")
    d.metric_card(sl, 0.5, 1.8, 3.5, 1.5, "月總收入", "206,859", GOLD, "薪資+房租+配息+股息")
    d.metric_card(sl, 4.3, 1.8, 3.5, 1.5, "月總支出", "210,958", RED, "房貸+利息+信用卡+生活")
    d.metric_card(sl, 8.1, 1.8, 3.5, 1.5, "月缺口", "-4,099 ⚠️", RED, "需仰賴配息填補")
    d.metric_card(sl, 0.5, 3.6, 3.5, 1.5, "保單借貸", "4,000,000", RED, "利率 ~5% 最高成本")
    d.metric_card(sl, 4.3, 3.6, 3.5, 1.5, "月保單利息", "-16,000", RED, "吃掉配息 18%")
    d.metric_card(sl, 8.1, 3.6, 3.5, 1.5, "若清償後月省", "+16,000 ✅", GREEN, "配息實收 73K→89K")
    d.card(sl, 0.5, 5.5, 11.5, 1.2)
    d.multi(sl, [
        "💡 轉貸資金成本僅 2.185%，而保單借貸利率 ~5%。每清償 100 萬保單借貸，年省利息 ~50,000。",
        "   清償 400 萬保單借貸 = 年省 192,000 利息 = 每月多 16,000 現金流。"
    ], 0.8, 5.7, 11, 0.8, 14, GOLD)

    # S4: 轉貸方案
    sl = d._new()
    d.page_title(sl, "築巢優利貸 2.185% 為最優解，月省 49,658", "三方案比較：現狀 vs 國泰轉貸 vs 築巢優利貸")
    d.metric_card(sl, 0.5, 1.8, 3.5, 1.8, "現狀：永豐房貸", "99,458/月", RED, "利率 ~2.5% ｜ 三筆分散")
    d.metric_card(sl, 4.3, 1.8, 3.5, 1.8, "方案 B：國泰轉貸", "52,500/月", GOLD, "利率 2.6% ｜ 次佳選擇")
    d.metric_card(sl, 8.1, 1.8, 3.5, 1.8, "方案 A：築巢優利貸 ⭐", "49,800/月", GREEN, "利率 2.185% ｜ 公務員專案")
    d.card(sl, 0.5, 4, 11.5, 2.7)
    d.txt(sl, "為什麼推薦築巢優利貸？", 0.8, 4.2, 8, 0.5, 20, True, GOLD)
    d.multi(sl, [
        "✅ 公務員身份符合，資格審查無虞",
        "✅ 2.185% 為市場最低房貸利率區間，鎖定固定利率不受升息影響",
        "✅ 月付從 99,458 降至 49,800，月釋放 49,658 現金流",
        "✅ 轉貸資金同時清償 4M 高利保單借貸，月再省 16,000",
        "⚠️ 9/25 現有房貸到期，需在此之前完成轉貸程序",
    ], 0.8, 4.8, 10, 1.5, 14)

    # S5: 三階段總覽
    sl = d._new()
    d.page_title(sl, "三階段資金部署：清償→安全網→低利擴張", "總資金 ~7M，目標月現金流改善 +106,358")
    phases = [
        ("第一階段", "清償高利負債", "4,000,000", "+16,000/月", GREEN, "安聯A 2M\n安聯B 1M\n第一金 1M"),
        ("第二階段", "建立安全網", "3,100,000", "備援無虞", BLUE, "保留理財型額度 3M\n星展備用金 100K"),
        ("第三階段", "低利擴張投資", "~3,000,000", "+7,200↑/月", GOLD, "ETF 分批建倉\n保單第三站 400 萬"),
    ]
    for i, (phase, title, amount, result, color, detail) in enumerate(phases):
        x = 0.5 + i * 4.2
        d.card(sl, x, 1.8, 3.8, 5)
        d.card(sl, x, 1.8, 3.8, 0.06)
        d.txt(sl, phase, x + 0.3, 2, 3, 0.3, 14, False, GRAY)
        d.txt(sl, title, x + 0.3, 2.4, 3, 0.4, 22, True, color)
        d.txt(sl, amount, x + 0.3, 3.1, 3, 0.5, 24, True, WHITE)
        d.txt(sl, result, x + 0.3, 3.7, 3, 0.3, 14, False, color)
        d.multi(sl, detail.split('\n'), x + 0.3, 4.4, 3, 1.5, 12, GRAY)

    # S6-S8: 三階段細節
    for phase_num, (title, action, items, summary) in enumerate([
        ("清償 4M 高利負債，月省 16,000", "保單配息實收從 73,000 提升至 89,000，年化節省 192,000", [
            ("安聯A 保單借貸", "2,000,000", "~5%", "~8,000", GREEN),
            ("安聯B 保單借貸", "1,000,000", "~5%", "~4,000", GREEN),
            ("第一金 FL65 保單借貸", "1,000,000", "~5%", "~4,000", GREEN),
        ], "✅ 清償後保單配息全額實領，無提前清償違約金"),
        ("保留 3M 額度 + 補足備用金 100K", "不動用不計息，保留財務韌性以因應黑天鵝", [
            ("保留理財型額度", "3,000,000", "備而不用", "黑天鵝緩衝", BLUE),
            ("星展活存補足", "100,000", "目前 17K→100K", "6個月生活費", BLUE),
        ], "✅ 理財型房貸特色：隨借隨還，按日計息，最適合緊急備援"),
    ]):
        sl = d._new()
        d.page_title(sl, title, action)
        for i, (name, amount, rate, saving, color) in enumerate(items):
            y = 2 + i * 1.3
            d.card(sl, 0.5, y, 12, 1)
            d.txt(sl, name, 0.8, y + 0.1, 4, 0.4, 18, True)
            d.txt(sl, amount, 5.5, y + 0.1, 2.5, 0.5, 26, True, WHITE)
            d.txt(sl, rate, 8.2, y + 0.1, 2, 0.4, 14, False, GRAY)
            d.txt(sl, saving, 10.5, y + 0.1, 2, 0.4, 18, True, color)
        d.card(sl, 0.5, 5.5, 12, 0.8)
        d.txt(sl, summary, 0.8, 5.65, 11, 0.5, 14, False, GREEN)

    # S9: 低利擴張
    sl = d._new()
    d.page_title(sl, "低利套利 2.185%，目標報酬率 >6%", "資金成本 2.185%  vs  目標資產報酬率 6-8%  →  利差 ~4%")
    for i, (title, amount, strategy, ret, color) in enumerate([
        ("台股 ETF 核心建倉", "~1,200,000", "逐月買入 00983D/00919/00878", "+7,200/月", GREEN),
        ("保單第三站配置", "4,000,000", "PIMCO 收益增長 + AI + A10", "+17,500/月", PURPLE),
        ("009816 成長型累積", "~400,000", "凱基台灣 TOP 50 不配息", "長期資本利得", BLUE),
    ]):
        y = 2.2 + i * 1.5
        d.card(sl, 0.5, y, 12, 1.2)
        d.card(sl, 0.5, y, 0.06, 1.2)
        d.txt(sl, title, 0.8, y + 0.1, 4, 0.4, 18, True, color)
        d.txt(sl, amount, 5.5, y + 0.1, 2, 0.5, 22, True, WHITE)
        d.txt(sl, strategy, 0.8, y + 0.65, 7, 0.4, 13, False, GRAY)
        d.txt(sl, ret, 9.5, y + 0.3, 2.5, 0.4, 18, True, color)

    # S10: 時間軸
    d.slide_timeline(
        title="7 月至 12 月執行路徑，9 月為關鍵轉折點",
        action_line="Q3 準備 → Q4 執行 → 年底檢視",
        events=[
            ("Q3 · 7月 ✅", "國泰轉貸面簽/對保完成", GREEN),
            ("Q3 · 8月", "確認轉貸細節 · 逐月買入 ETF", GRAY),
            ("Q3 · 9/25 🔴", "永豐房貸到期 · 國泰撥款 · 清償 4M 保單借貸", RED),
            ("Q4 · 10/1", "第二階段：安全網 + 辦理築巢優利貸 2.185%", BLUE),
            ("Q4 · 10-12月", "第三階段：ETF/保單第三站佈局", GOLD),
            ("Q4 · 12月底", "全年檢視成效 · CIO 審查", GRAY),
        ])

    # S11: 成效對照
    sl = d._new()
    d.page_title(sl, "轉貸後月淨現金流從 -4,099 提升至 +102,259", "負債率從 41.8% 降至 ~35%，月現金流改善 +106,358")
    cd = CategoryChartData()
    cd.categories = ['轉貸前', '轉貸後']
    cd.add_series('月收入', (206.9, 247.6))
    cd.add_series('月支出', (211.0, 145.3))
    sl.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, Inches(0.5), Inches(1.8), Inches(6.5), Inches(4.5), cd)
    metrics = [
        ('負債率', '41.8% → ~35%', GREEN),
        ('房貸月付', '99,458 → 49,800', GREEN),
        ('保單利息', '16,000 → 0', GREEN),
        ('月配息實收', '73,000 → 106,500', GREEN),
        ('月淨現金流', '-4,099 → +102,259', GOLD),
    ]
    for i, (label, value, color) in enumerate(metrics):
        y = 2 + i * 0.85
        d.card(sl, 7.5, y, 5, 0.65)
        d.txt(sl, label, 7.8, y + 0.08, 2, 0.3, 13, False, GRAY)
        d.txt(sl, value, 10, y + 0.1, 2.5, 0.4, 17, True, color)

    # S12: 風險
    d.slide_risks(
        title="利率、審核、市場風險皆可控，執行風險低",
        action_line="每項風險都有對應緩解措施，最壞情境不影響核心財務",
        risks=[
            ("利率上升風險", "築巢優利貸若隨央行升息調整", "固定利率已鎖 2.185%"),
            ("轉貸審核未過", "無法取得低利資金", "國泰已面簽/對保 ✅"),
            ("房價下跌", "LTV 不足須補擔保", "負債率 41.8%，安全邊際充足"),
            ("ETF 價格下跌", "投入本金虧損", "分批買入 + 配息保護"),
            ("保單配息縮水", "基金淨值影響配息", "分散標的 + 核心衛星配置"),
        ])

    # S13: 總結
    d.slide_summary("結論：每月從 -4,099 到 +102,259", [
        ("🎯", "清償高利", "消滅 4M 保單借貸（5%），年省 192K 利息，無風險套利", GREEN),
        ("🛡️", "保留韌性", "3M 理財型額度 + 100K 備用金，財務緩衝充足", BLUE),
        ("📈", "低利套利", "2.185% 資金投入 >6% 資產，賺取 ~4% 利差", GOLD),
        ("💡", "現金流翻轉", "月淨現金流從 -4,099 提升至 +102,259，財務自由加速", GREEN),
    ])

    path = d.save("C:/Users/bot/Desktop/longjiu_system/轉貸投資簡報_v3.pptx")
    print(f"✅ {path} (13 頁)")
