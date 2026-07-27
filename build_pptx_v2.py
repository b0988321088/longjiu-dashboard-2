"""轉貸投資簡報 v2 — 15頁，深色主題，高對比"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.chart import XL_CHART_TYPE, XL_LABEL_POSITION
from pptx.dml.color import RGBColor
from pptx.chart.data import CategoryChartData

prs = Presentation()
prs.slide_width = Inches(13.33)
prs.slide_height = Inches(7.5)

# 品牌色系 — 高對比
BG = RGBColor(0x0B, 0x0D, 0x1A)        # 深藍黑底
CARD = RGBColor(0x15, 0x18, 0x2E)       # 卡片底色
CARD_BORDER = RGBColor(0x1E, 0x22, 0x3D) # 卡片邊框
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
TEXT = RGBColor(0xE8, 0xE8, 0xF0)       # 主文字
GRAY = RGBColor(0x8A, 0x8F, 0xA0)       # 輔助文字
GOLD = RGBColor(0xF7, 0xA0, 0x1C)       # 強調金
GREEN = RGBColor(0x34, 0xD3, 0x99)      # 正
RED = RGBColor(0xFF, 0x5C, 0x5C)        # 負
BLUE = RGBColor(0x5B, 0x9B, 0xF7)       # 資訊
PURPLE = RGBColor(0xA7, 0x8B, 0xFA)     # 次要

def new_slide():
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    sl.background.fill.solid()
    sl.background.fill.fore_color.rgb = BG
    return sl

def add_shape(sl, left, top, w, h, fill=CARD, line=CARD_BORDER):
    s = sl.shapes.add_shape(1, Inches(left), Inches(top), Inches(w), Inches(h))
    s.fill.solid(); s.fill.fore_color.rgb = fill
    s.line.color.rgb = line; s.line.width = Pt(1)
    return s

def txt(sl, text, left, top, w, h, size=14, bold=False, color=TEXT, align=PP_ALIGN.LEFT):
    tb = sl.shapes.add_textbox(Inches(left), Inches(top), Inches(w), Inches(h))
    tb.text_frame.word_wrap = True
    p = tb.text_frame.paragraphs[0]
    p.text = text; p.font.size = Pt(size); p.font.bold = bold
    p.font.color.rgb = color; p.alignment = align
    return tb

def multi_line(sl, lines, left, top, w, h, size=13, color=TEXT):
    """多行文字，每行用 \n 分隔"""
    tb = sl.shapes.add_textbox(Inches(left), Inches(top), Inches(w), Inches(h))
    tb.text_frame.word_wrap = True
    for i, line in enumerate(lines):
        if i == 0:
            p = tb.text_frame.paragraphs[0]
        else:
            p = tb.text_frame.add_paragraph()
        p.text = line; p.font.size = Pt(size); p.font.color.rgb = color
        p.space_after = Pt(4)

# ============ SLIDE 1: 封面 ============
sl = new_slide()
txt(sl, '龍 九 控 股', 1, 1, 5, 1, 16, False, GRAY)
txt(sl, '轉貸投資　可行性評估計劃', 1, 1.8, 11, 1.5, 44, True, WHITE)
# 分隔線
add_shape(sl, 1, 3.5, 4, 0.04, fill=GOLD)
txt(sl, '以低利資金優化負債結構 × 提升被動現金流', 1, 3.9, 11, 0.6, 20, False, GRAY)
txt(sl, '2026 年 7 月 26 日 ｜ 評估期間：2026 Q3–Q4 ｜ CIO 策略研究室', 1, 5.8, 8, 0.5, 14, False, GRAY)

# ============ SLIDE 2: 目錄 ============
sl = new_slide()
txt(sl, '目錄 CONTENTS', 1, 0.5, 6, 0.8, 32, True, WHITE)
add_shape(sl, 1, 1.3, 3, 0.03, fill=GOLD)
items = [
    '01 · 現狀診斷：資產負債結構',
    '02 · 現狀診斷：每月收支缺口',
    '03 · 轉貸方案評估與比較',
    '04 · 三階段資金部署計劃',
    '05 · 第一階段：清償高利負債',
    '06 · 第二階段：建立安全網',
    '07 · 第三階段：低利擴張投資',
    '08 · ETF 建倉策略',
    '09 · 保單第三站配置',
    '10 · 執行時間表',
    '11 · 預期成效對照',
    '12 · 敏感度分析與風險',
    '13 · 總結與建議',
]
for i, item in enumerate(items):
    y = 2 + i * 0.4
    txt(sl, item, 1.5, y, 10, 0.35, 15, color=TEXT)

# ============ SLIDE 3: 資產負債結構 ============
sl = new_slide()
txt(sl, '01  現狀診斷：資產負債結構', 0.8, 0.4, 10, 0.6, 26, True, WHITE)
add_shape(sl, 0.8, 1.1, 3, 0.03, fill=GOLD)
txt(sl, f'總資產 50,689,930  ｜  總負債 21,165,869  ｜  負債率 41.8%', 0.8, 1.4, 10, 0.4, 14, False, GRAY)

cd = CategoryChartData()
cd.categories = ['不動產\n33.3M','保單\n9.8M','證券\n2.5M','現金\n3.3M','基金\n0.8M']
cd.add_series('', (33.3, 9.8, 2.5, 3.3, 0.8))
c = sl.shapes.add_chart(XL_CHART_TYPE.DOUGHNUT, Inches(0.5), Inches(2), Inches(5.5), Inches(4.5), cd).chart
c.has_legend = True

cd2 = CategoryChartData()
cd2.categories = ['房貸\n13.2M','理財型\n3.0M','保單借貸\n4.0M','質押\n1.0M']
cd2.add_series('', (13.2, 3.0, 4.0, 1.0))
c2 = sl.shapes.add_chart(XL_CHART_TYPE.DOUGHNUT, Inches(6.5), Inches(2), Inches(5.5), Inches(4.5), cd2).chart
c2.has_legend = True

add_shape(sl, 0.5, 6.5, 5.5, 0.5)
multi_line(sl, ['■ 資產以不動產（65.7%）為主，流動性偏低', '■ 保單借貸 4M 利率 ~5%，為最高成本負債'], 0.8, 6.6, 5, 0.4, 11, GRAY)
add_shape(sl, 6.5, 6.5, 5.5, 0.5)
multi_line(sl, ['■ 三筆房貸利率分散，整合後可降低月付', '■ 證券質押 1M 因 0056 凍結暫時無法操作'], 6.8, 6.6, 5, 0.4, 11, GRAY)

# ============ SLIDE 4: 每月收支 ============
sl = new_slide()
txt(sl, '02  現狀診斷：每月收支缺口', 0.8, 0.4, 10, 0.6, 26, True, WHITE)
add_shape(sl, 0.8, 1.1, 3, 0.03, fill=GOLD)
txt(sl, '月收入 ~206,859  ｜  月支出 ~210,958  ｜  月缺口 -4,099 ⚠️', 0.8, 1.4, 12, 0.4, 14, False, GRAY)

cd3 = CategoryChartData()
cd3.categories = ['薪資','房租','保單配息\n(扣息後)','股息','基金']
cd3.add_series('收入 (K)', (43.1, 80.1, 73.0, 10.0, 0.6))
c3 = sl.shapes.add_chart(XL_CHART_TYPE.BAR_CLUSTERED, Inches(0.5), Inches(2), Inches(5.5), Inches(4.5), cd3).chart

cd4 = CategoryChartData()
cd4.categories = ['房貸\n99.5K','理財利息\n10K','保單利息\n16K','信用卡\n38K','生活\n45K','質押\n2.5K']
cd4.add_series('支出 (K)', (99.5, 10.0, 16.0, 38.0, 45.0, 2.5))
c4 = sl.shapes.add_chart(XL_CHART_TYPE.BAR_CLUSTERED, Inches(6.5), Inches(2), Inches(5.5), Inches(4.5), cd4).chart

add_shape(sl, 3.5, 6.7, 5, 0.5, fill=RGBColor(0x2D, 0x1B, 0x1B))
txt(sl, '⚠️ 每月入不敷出，需仰賴保單配息填補缺口', 3.8, 6.75, 5, 0.4, 13, True, RED)

# ============ SLIDE 5: 轉貸方案 ============
sl = new_slide()
txt(sl, '03  轉貸方案評估與比較', 0.8, 0.4, 10, 0.6, 26, True, WHITE)
add_shape(sl, 0.8, 1.1, 3, 0.03, fill=GOLD)

cards_data = [
    ('方案A：築巢優利貸', '2.185%', '49,800/月', '✅ 最優', GOLD),
    ('方案B：國泰轉貸', '2.6%', '52,500/月', '⭐ 次佳', BLUE),
    ('方案C：維持現狀', '~2.5%', '99,458/月', '❌ 不佳', RED),
]
for i, (title, rate, monthly, rec, color) in enumerate(cards_data):
    x = 0.5 + i * 4.2
    add_shape(sl, x, 2, 3.8, 3.5, fill=CARD)
    add_shape(sl, x, 2, 3.8, 0.06, fill=color)
    txt(sl, title, x + 0.3, 2.3, 3.2, 0.4, 18, True, color)
    txt(sl, f'利率 {rate}', x + 0.3, 3, 3.2, 0.4, 28, True, WHITE)
    txt(sl, f'月付 {monthly}', x + 0.3, 3.7, 3.2, 0.4, 16, False, GRAY)
    txt(sl, rec, x + 0.3, 4.5, 3.2, 0.4, 16, True, color)

add_shape(sl, 0.5, 6, 12, 1)
multi_line(sl, [
    '推薦方案 A：築巢優利貸 2.185% ｜ 公務員專案，資格符合',
    '每月房貸從 99,458 降至 49,800，月省 49,658 ｜ 轉貸資金可同步清償 4M 保單借貸',
], 0.8, 6.15, 11, 0.8, 14, GRAY)

# ============ SLIDE 6: 三階段總覽 ============
sl = new_slide()
txt(sl, '04  三階段資金部署計劃', 0.8, 0.4, 10, 0.6, 26, True, WHITE)
add_shape(sl, 0.8, 1.1, 3, 0.03, fill=GOLD)

phases = [
    ('第一階段', '清償高利負債', '4,000,000', '月省 16,000', GREEN, '清償安聯A 2M\n安聯B 1M\n第一金 FL65 1M'),
    ('第二階段', '建立安全網', '3,100,000', '備援無虞', BLUE, '保留理財型額度 3M\n補足星展備用金 100K'),
    ('第三階段', '低利擴張投資', '~3,000,000', '月增 ~7,000↑', GOLD, 'ETF 分批建倉\n保單第三站 400 萬'),
]
for i, (phase, title, amount, result, color, detail) in enumerate(phases):
    x = 0.5 + i * 4.2
    add_shape(sl, x, 1.8, 3.8, 5)
    add_shape(sl, x, 1.8, 3.8, 0.06, fill=color)
    txt(sl, phase, x + 0.3, 2, 3.2, 0.4, 14, False, GRAY)
    txt(sl, title, x + 0.3, 2.4, 3.2, 0.4, 20, True, color)
    txt(sl, f'資金 {amount}', x + 0.3, 3.1, 3.2, 0.4, 22, True, WHITE)
    txt(sl, result, x + 0.3, 3.7, 3.2, 0.4, 14, False, color)
    multi_line(sl, detail.split('\n'), x + 0.3, 4.4, 3, 1.5, 12, GRAY)

# ============ SLIDE 7: 清償高利 ============
sl = new_slide()
txt(sl, '05  第一階段：清償高利負債', 0.8, 0.4, 10, 0.6, 26, True, WHITE)
add_shape(sl, 0.8, 1.1, 3, 0.03, fill=GREEN)
txt(sl, '撥款後 D+1 立即執行  ｜  消除年化 5% 利息支出', 0.8, 1.4, 10, 0.4, 14, False, GRAY)

items_data = [
    ('安聯A 保單借貸', '2,000,000', '~5%', '~8,000', GREEN),
    ('安聯B 保單借貸', '1,000,000', '~5%', '~4,000', GREEN),
    ('第一金 FL65 保單借貸', '1,000,000', '~5%', '~4,000', GREEN),
]
for i, (name, amount, rate, saving, color) in enumerate(items_data):
    y = 2.2 + i * 1.3
    add_shape(sl, 0.5, y, 12, 1)
    txt(sl, name, 0.8, y + 0.1, 4, 0.4, 16, True)
    txt(sl, amount, 5.5, y + 0.1, 2.5, 0.4, 24, True, WHITE)
    txt(sl, f'利率 {rate}', 8, y + 0.1, 2, 0.4, 14, False, RED)
    txt(sl, f'月省 {saving}', 10, y + 0.1, 2, 0.4, 18, True, GREEN)

add_shape(sl, 0.5, 6.2, 12, 0.8, fill=RGBColor(0x0D, 0x2B, 0x1E))
multi_line(sl, [
    '✅ 清償後保單配息全額實領：73,000 → 89,000（+16,000/月）',
    '✅ 年化節省利息：192,000  ｜  無提前清償違約金',
], 0.8, 6.35, 11, 0.7, 13, GREEN)

# ============ SLIDE 8: 安全網 ============
sl = new_slide()
txt(sl, '06  第二階段：建立安全網', 0.8, 0.4, 10, 0.6, 26, True, WHITE)
add_shape(sl, 0.8, 1.1, 3, 0.03, fill=BLUE)

nets = [
    ('保留理財型房貸額度', '3,000,000', '不關閉，備而不用', '黑天鵝緩衝'),
    ('星展活存補足', '100,000', '目前 17,000 → 補至 100K', '6 個月生活費'),
]
for i, (title, amount, detail, purpose) in enumerate(nets):
    y = 2 + i * 1.8
    add_shape(sl, 0.5, y, 12, 1.5)
    txt(sl, title, 0.8, y + 0.15, 4, 0.4, 20, True, BLUE)
    txt(sl, amount, 5.5, y + 0.15, 2.5, 0.5, 32, True, WHITE)
    txt(sl, detail, 0.8, y + 0.7, 6, 0.4, 13, False, GRAY)
    txt(sl, purpose, 8, y + 0.7, 4, 0.4, 13, False, GOLD)

add_shape(sl, 0.5, 5.8, 12, 1)
multi_line(sl, [
    '安全網設計原則：不動用不計息，保留財務韌性',
    '理財型房貸特色：隨借隨還，按日計息，最適合緊急備援',
], 0.8, 5.95, 11, 0.8, 13, GRAY)

# ============ SLIDE 9: 低利擴張 ============
sl = new_slide()
txt(sl, '07  第三階段：低利擴張投資', 0.8, 0.4, 10, 0.6, 26, True, WHITE)
add_shape(sl, 0.8, 1.1, 3, 0.03, fill=GOLD)
txt(sl, f'資金成本 2.185%  ｜  目標報酬率 > 6%  ｜  利差 ~4%', 0.8, 1.4, 10, 0.4, 14, False, GRAY)

investments = [
    ('台股 ETF 核心建倉', '~1,200,000', '逐月買入 00983D/00919/00878', '預估 +7,200/月', GREEN),
    ('保單第三站配置', '4,000,000', 'PIMCO 收益增長 + AI + A10', '預估 +17,500/月', PURPLE),
    ('009816 成長型累積', '~400,000', '凱基台灣 TOP 50 不配息', '長期資本利得', BLUE),
]
for i, (title, amount, strategy, return_val, color) in enumerate(investments):
    y = 2.2 + i * 1.5
    add_shape(sl, 0.5, y, 12, 1.2)
    add_shape(sl, 0.5, y, 0.06, 1.2, fill=color)
    txt(sl, title, 0.8, y + 0.1, 4, 0.4, 18, True, color)
    txt(sl, amount, 5.5, y + 0.1, 2, 0.4, 20, True, WHITE)
    txt(sl, strategy, 0.8, y + 0.6, 7, 0.4, 13, False, GRAY)
    txt(sl, return_val, 9, y + 0.3, 3, 0.4, 16, True, color)

# ============ SLIDE 10: ETF 建倉 ============
sl = new_slide()
txt(sl, '08  ETF 建倉策略', 0.8, 0.4, 10, 0.6, 26, True, WHITE)
add_shape(sl, 0.8, 1.1, 3, 0.03, fill=GREEN)

etfs = [
    ('00983D', '主動富邦複合收益', '10→100張', '每月除息後買 10 張', '月配/債+股', GREEN),
    ('00919', '群益精選高息', '補至 10 張', '除息後補滿', '季配/高股息', GREEN),
    ('00878', '國泰永續高股息', '補至 20 張', '除息後補滿', '季配/高股息', GREEN),
    ('009816', '凱基台灣TOP50', '16→30張', '每月累積', '不配息/成長', BLUE),
    ('0050/006208', '台灣50', '各 2→5 張', '季線以下買', '半年配/市值', PURPLE),
]
for i, (code, name, target, strategy, note, color) in enumerate(etfs):
    y = 1.8 + i * 1
    add_shape(sl, 0.5, y, 12, 0.8)
    txt(sl, code, 0.8, y + 0.1, 1.5, 0.4, 18, True, color)
    txt(sl, name, 2.5, y + 0.1, 3, 0.4, 14, False)
    txt(sl, target, 5.5, y + 0.1, 1.5, 0.4, 16, True, WHITE)
    txt(sl, strategy, 7.5, y + 0.1, 3.5, 0.4, 12, False, GRAY)
    txt(sl, note, 11, y + 0.1, 2, 0.4, 11, False, color)

# ============ SLIDE 11: 保單第三站 ============
sl = new_slide()
txt(sl, '09  保單第三站配置', 0.8, 0.4, 10, 0.6, 26, True, WHITE)
add_shape(sl, 0.8, 1.1, 3, 0.03, fill=PURPLE)
txt(sl, '總資金 4,000,000  ｜  資產配置：47% 債券 + 53% 科技', 0.8, 1.4, 10, 0.4, 14, False, GRAY)

funds = [
    ('PIMCO 收益增長 M 級', '47%', '1,880,000', '~8,000/月', '債券為主 + 股票增益'),
    ('安聯 AI 收益成長 B 型', '30%', '1,200,000', '~5,500/月', '多重資產 + AI 主題'),
    ('貝萊德世界科技 A10', '23%', '920,000', '~4,000/月', '科技股 + 月配息'),
]
for i, (name, pct, amount, est, note) in enumerate(funds):
    y = 2.2 + i * 1.4
    add_shape(sl, 0.5, y, 12, 1.1)
    txt(sl, name, 0.8, y + 0.1, 4, 0.4, 16, True)
    txt(sl, pct, 5, y + 0.1, 1, 0.4, 20, True, PURPLE)
    txt(sl, amount, 6.5, y + 0.1, 2, 0.4, 18, True, WHITE)
    txt(sl, est, 9, y + 0.1, 2.5, 0.4, 16, True, GREEN)
    txt(sl, note, 0.8, y + 0.6, 8, 0.4, 12, False, GRAY)

add_shape(sl, 0.5, 6.2, 12, 0.7)
multi_line(sl, [
    '預計月配息合計 ~17,500  ｜  註：第三站已在第一金 FL65 轉換中，7 月底入帳',
], 0.8, 6.35, 11, 0.5, 13, GRAY)

# ============ SLIDE 12: 時間表 ============
sl = new_slide()
txt(sl, '10  執行時間表', 0.8, 0.4, 10, 0.6, 26, True, WHITE)
add_shape(sl, 0.8, 1.1, 3, 0.03, fill=GOLD)

timeline = [
    ('Q3 7月 ✅', '國泰轉貸面簽/對保完成', GREEN),
    ('Q3 8月', '確認轉貸細節 · 逐月買入 ETF', GRAY),
    ('Q3 9/25 🔴', '永豐房貸到期 · 國泰撥款', RED),
    ('Q3 9月底', '第一階段：清償 4M 保單借貸', GREEN),
    ('Q4 10/1', '第二階段：安全網 + 辦理築巢優利貸', BLUE),
    ('Q4 10-12月', '第三階段：ETF/保單佈局', GOLD),
    ('Q4 12月底', '全年檢視成效 · CIO 審查', GRAY),
]
for i, (date, event, color) in enumerate(timeline):
    y = 1.8 + i * 0.75
    add_shape(sl, 0.5, y, 12, 0.6)
    txt(sl, date, 0.8, y + 0.1, 2, 0.4, 14, True, color)
    txt(sl, event, 3.5, y + 0.1, 8, 0.4, 14, False)

# ============ SLIDE 13: 成效 ============
sl = new_slide()
txt(sl, '11  預期成效對照', 0.8, 0.4, 10, 0.6, 26, True, WHITE)
add_shape(sl, 0.8, 1.1, 3, 0.03, fill=GREEN)

cd5 = CategoryChartData()
cd5.categories = ['轉貸前', '轉貸後']
cd5.add_series('月收入 (K)', (206.9, 247.6))
cd5.add_series('月支出 (K)', (211.0, 145.3))
c5 = sl.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, Inches(0.5), Inches(1.8), Inches(6.5), Inches(4.5), cd5).chart

metrics = [
    ('負債率', '41.8% → ~35%', GREEN),
    ('房貸月付', '99,458 → 49,800', GREEN),
    ('保單利息', '16,000 → 0', GREEN),
    ('月配息實收', '73,000 → 106,500', GREEN),
    ('月淨現金流', '-4,099 → +102,259', GOLD),
]
for i, (label, value, color) in enumerate(metrics):
    y = 2 + i * 0.8
    add_shape(sl, 7.5, y, 5, 0.6)
    txt(sl, label, 7.8, y + 0.1, 2, 0.4, 13, False, GRAY)
    txt(sl, value, 10, y + 0.1, 2.5, 0.4, 16, True, color)

# ============ SLIDE 14: 風險 ============
sl = new_slide()
txt(sl, '12  敏感度分析與風險', 0.8, 0.4, 10, 0.6, 26, True, WHITE)
add_shape(sl, 0.8, 1.1, 3, 0.03, fill=RED)

risks = [
    ('利率上升風險', '築巢優利貸若隨央行升息調整', '固定利率已鎖 2.185%，影響有限'),
    ('轉貸審核未過', '無法取得低利資金', '國泰已面簽/對保 ✅ 風險極低'),
    ('房價下跌', 'LTV 不足需補擔保', '負債率 41.8%，安全邊際充足'),
    ('ETF 價格下跌', '投入本金虧損', '分批買入 + 配息保護 + 長期持有'),
    ('保單配息縮水', '基金淨值波動影響配息', '分散標的 + 核心衛星配置'),
]
for i, (risk, impact, mitigation) in enumerate(risks):
    y = 1.8 + i * 1
    add_shape(sl, 0.5, y, 12, 0.8)
    txt(sl, risk, 0.8, y + 0.1, 3, 0.4, 15, True, RED if '風險' in risk else GOLD)
    txt(sl, impact, 4, y + 0.1, 4, 0.4, 12, False, GRAY)
    txt(sl, '✅ ' + mitigation, 8.5, y + 0.1, 4, 0.4, 12, False, GREEN)

# ============ SLIDE 15: 總結 ============
sl = new_slide()
txt(sl, '13  總結與建議', 0.8, 0.4, 10, 0.6, 26, True, WHITE)
add_shape(sl, 0.8, 1.1, 3, 0.03, fill=GOLD)

conclusion = [
    ('🎯', '清償高利', '優先消滅 4M 保單借貸（5%），年省 192K 利息', GREEN),
    ('🛡️', '保留韌性', '保留 3M 理財型額度 + 補足備用金 100K', BLUE),
    ('📈', '低利套利', '2.185% 資金投入 >6% 收益資產，賺取 ~4% 利差', GOLD),
    ('📊', '分批建倉', 'ETF 逐月買入 + 保單第三站 400 萬，分散風險', PURPLE),
    ('💡', '現金流翻轉', '月淨現金流從 -4,099 提升至 +102,259', GREEN),
]
for i, (icon, title, desc, color) in enumerate(conclusion):
    y = 1.5 + i * 1
    add_shape(sl, 0.5, y, 12, 0.8)
    txt(sl, f'{icon}  {title}', 0.8, y + 0.1, 3, 0.4, 20, True, color)
    txt(sl, desc, 4, y + 0.15, 8, 0.4, 14, False, TEXT)

txt(sl, '本計劃由龍九控股 CIO 策略研究室製作  ｜  數據來源：snapshot.json / dragon_assets.db', 1, 7, 12, 0.4, 11, False, GRAY)

prs.save('C:/Users/bot/Desktop/longjiu_system/轉貸投資簡報_v2.pptx')
print('✅ 轉貸投資簡報_v2.pptx (15頁)')
