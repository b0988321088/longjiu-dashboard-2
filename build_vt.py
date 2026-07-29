"""
新增 VT vs 0050 分析頁
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
import json, os

BASE = os.path.dirname(os.path.abspath(__file__))
SNAP = json.load(open(f'{BASE}/snapshot.json'))
INS = SNAP.get('allianz_combined',0) + SNAP.get('firstjin_fl65_current_value',0)
SEC = SNAP.get('securities_total_market_value',0)
FUND = SNAP.get('fund_market_value',0)
CASH = SNAP.get('real_liquid_assets',0)
p = SNAP.get('penetration',{}).get('actual_pct',{})

prs = Presentation()
prs.slide_width = Inches(13.33); prs.slide_height = Inches(7.5)
BG = RGBColor(0x0B,0x0D,0x1A); WHITE = RGBColor(0xFF,0xFF,0xFF)
GRAY = RGBColor(0x8A,0x8F,0xA0); GOLD = RGBColor(0xF7,0xA0,0x1C)
GREEN = RGBColor(0x34,0xD3,0x99); RED = RGBColor(0xFF,0x5C,0x5C)

def ns():
    s = prs.slides.add_slide(prs.slide_layouts[6])
    s.background.fill.solid(); s.background.fill.fore_color.rgb = BG
    return s

def T(s, t, top=0.3):
    tb = s.shapes.add_textbox(Inches(0.8), Inches(top), Inches(11.5), Inches(0.8))
    p = tb.text_frame.paragraphs[0]
    p.text = t; p.font.size = Pt(32); p.font.bold = True; p.font.color.rgb = WHITE

def ST(s, t, top=1.2):
    tb = s.shapes.add_textbox(Inches(0.8), Inches(top), Inches(11.5), Inches(0.5))
    p = tb.text_frame.paragraphs[0]
    p.text = t; p.font.size = Pt(16); p.font.color.rgb = GRAY

def B(s, items, top=1.9):
    tb = s.shapes.add_textbox(Inches(1), Inches(top), Inches(11.3), Inches(5))
    tf = tb.text_frame; tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i==0 else tf.add_paragraph()
        p.text = item; p.font.size = Pt(15); p.font.color.rgb = WHITE
        p.space_after = Pt(5)

# === VT vs 0050 分析 ===
s = ns()
T(s, 'VT vs 0050：該買哪個？')
ST(s, '核心問題：您需要的是全球分散，還是台灣集中？')

B(s, [
    '🌍 VT（Vanguard 全球股票 ETF）：',
    f'     60% 美股 + 40% 國際（歐/日/新興），配息 ~2%',
    '      → 您已有安聯全球多資產（776 萬），功能重疊',
    '      → 需換美元（台幣 32.38 偏貴），且美股 33.7% 已超標',
    '',
    '🏢 0050（元大台灣 50）：',
    f'     100% 台灣 50 大權值股，配息 ~3-4%',
    f'      → 台股目前 {p.get("台股市值型成長",0):.1f}%（目標 35%），缺口 20pp',
    '      → 用台幣買，不用換匯，今天崩跌更便宜',
    '',
    '📊 比較結論：',
    '     006208/0050 是台股市值型，VT 是全球分散型',
    '     您的配置缺的是「台股集中」不是「全球分散」',
    '     安聯收益增長已在做全球分散的工作',
    '',
    '🎯 建議：買 0050/006208 補台股缺口，不買 VT',
    '     若真要全球分散 → 加碼 00983D（全球複合債 7.1%，台幣計價）'
])

prs.save(f'{BASE}/VT_vs_0050.pptx')
print('✅ VT vs 0050 分析，1 頁')
