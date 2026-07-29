"""
update_pptx.py — 以模板為基礎，只取代文字內容
不修改任何版面/位置/字型/顏色
"""
from pptx import Presentation
import re

SRC = '龍九簡報模板.pptx'
DST = '轉貸投資計劃_最新.pptx'

# === 文字取代規則 ===
# key = 舊文字（模板中原有文字）
# value = 新文字（要取代成的內容）
repl = {
    # Slide 1: 封面 - 更新日期
    '2026 年 7 月 26 日 ｜ 評估期間：2026 Q3–Q4': '2026 年 7 月 28 日 ｜ 更新版 ｜ 評估期間：2026 Q3–Q4',
    
    # Slide 2: 目錄 - 更新總結數字
    '+106K/月': '+61K/月',
    '+102,259': '+61,559',
    '每月現金流缺口 -4,099，保單利息吃掉 16,000': '每月現金流缺口 -4,099 → 轉貸後 +61,559',

    # Slide 3: 核心問題 - 更新資產數字（在文字段落中）
    '保單借貸 4,000,000': '保單借貸 4,000,000',
    
    # Slide 5: 三階段部署 - 更新金額
    '目標月現金流改善 +106,358': '目標月現金流改善 +61,559',
    '+16,000/月': '+16,000/月 ✅',
    '+7,200↑/月': '+7,200↑/月 ✅',

    # Slide 8: 低利套利 - 更新配置
    '~1,200,000': '~1,000,000',
    '00983D/00919/00878': '00983D 分批 + 00919 補倉 + 00878 補倉',
    '+7,200/月': '+7,200/月（估）',

    # Slide 9: 時間軸 - 更新狀態
    '國泰轉貸面簽/對保完成': '國泰轉貸面簽/對保完成 ✅',

    # Slide 10: 成效 - 更新數字
    '73,000 → 106,500': '73,000 → 89,000',
    '-4,099 → +102,259': '-4,099 → +61,559',
    '99,458 → 49,800': '99,458 → 49,800 ✅',

    # Slide 12: 結論
    '+102,259': '+61,559',
}

print('🔍 開啟模板...')
prs = Presentation(SRC)
count = 0

for si, slide in enumerate(prs.slides):
    for shape in slide.shapes:
        if shape.has_text_frame:
            tf = shape.text_frame
            for para in tf.paragraphs:
                for run in para.runs:
                    orig = run.text
                    new = orig
                    for old_t, new_t in repl.items():
                        if old_t in new:
                            new = new.replace(old_t, new_t)
                    if new != orig:
                        run.text = new
                        count += 1

prs.save(DST)
print(f'✅ 完成！{count} 處文字已更新')
print(f'📁 {DST}')
