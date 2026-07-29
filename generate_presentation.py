
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import MSO_ANCHOR, MSO_AUTO_SIZE
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

# Presentation settings
prs = Presentation()
prs.slide_width = Inches(10)
prs.slide_height = Inches(7.5)

# Define custom colors based on the skill
# Main Background
BG = RGBColor(0x0B, 0x0D, 0x1A)
# Card Background
CARD = RGBColor(0x15, 0x18, 0x2E)
# Card Border (slightly lighter than CARD)
CARD_BORDER = RGBColor(0x1E, 0x22, 0x3D)
# Text Colors
TEXT_LIGHT = RGBColor(0xE8, 0xE8, 0xF0)
TEXT_GRAY = RGBColor(0x8A, 0x8F, 0xA0)
# Accent Colors
GREEN = RGBColor(0x34, 0xD3, 0x99)  # Positive/Gain
RED = RGBColor(0xFF, 0x5C, 0x5C)    # Warning/Negative
GOLD = RGBColor(0xF7, 0xA0, 0x1C)   # Highlight
BLUE = RGBColor(0x5B, 0x9B, 0xF7)   # Information

# Font settings
FONT_TITLE_SIZE = Pt(28)
FONT_SUBTITLE_SIZE = Pt(16)
FONT_CARD_TITLE_SIZE = Pt(18)
FONT_BODY_SIZE = Pt(13)
FONT_KEY_METRIC_SIZE = Pt(40)

# Helper function to add a title slide
def add_title_slide(prs, title_text, subtitle_text, date_text):
    slide_layout = prs.slide_layouts[0] # Title slide layout
    slide = prs.slides.add_slide(slide_layout)
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = BG

    title = slide.shapes.title
    title.text = title_text
    title.text_frame.paragraphs[0].font.color.rgb = TEXT_LIGHT
    title.text_frame.paragraphs[0].font.size = FONT_TITLE_SIZE
    title.text_frame.paragraphs[0].font.bold = True

    # Subtitle
    left = Inches(1)
    top = Inches(2.5)
    width = Inches(8)
    height = Inches(1)
    subtitle_box = slide.shapes.add_textbox(left, top, width, height)
    subtitle_tf = subtitle_box.text_frame
    p = subtitle_tf.add_paragraph()
    p.text = subtitle_text
    p.font.color.rgb = TEXT_GRAY
    p.font.size = FONT_SUBTITLE_SIZE

    # Date
    date_box = slide.shapes.add_textbox(Inches(1), Inches(6.5), Inches(8), Inches(0.5))
    date_tf = date_box.text_frame
    p_date = date_tf.add_paragraph()
    p_date.text = date_text
    p_date.font.color.rgb = TEXT_GRAY
    p_date.font.size = Pt(10)
    p_date.alignment = MSO_ANCHOR.BOTTOM

    return slide

# Helper function to add a standard content slide with dark background
def add_content_slide(prs, action_title, subtitle_insight):
    slide_layout = prs.slide_layouts[5] # Title and Content layout
    slide = prs.slides.add_slide(slide_layout)
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = BG

    # Action Title
    title = slide.shapes.title
    title.text = action_title
    title.text_frame.paragraphs[0].font.color.rgb = TEXT_LIGHT
    title.text_frame.paragraphs[0].font.size = FONT_TITLE_SIZE
    title.text_frame.paragraphs[0].font.bold = True

    # Subtitle/Insight (positioned below title)
    left = Inches(1)
    top = Inches(1.5) # Adjust position to be below the main title
    width = Inches(8)
    height = Inches(0.5)
    subtitle_box = slide.shapes.add_textbox(left, top, width, height)
    subtitle_tf = subtitle_box.text_frame
    p = subtitle_tf.add_paragraph()
    p.text = subtitle_insight
    p.font.color.rgb = TEXT_GRAY
    p.font.size = FONT_SUBTITLE_SIZE
    subtitle_tf.margin_bottom = Inches(0.1)
    subtitle_tf.margin_top = Inches(0.1)
    subtitle_tf.margin_left = Inches(0.1)
    subtitle_tf.margin_right = Inches(0.1)

    return slide

# Helper function to add a card
def add_card(slide, left, top, w, h):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = CARD
    shape.line.color.rgb = CARD_BORDER
    shape.line.width = Pt(1)
    return shape

# Helper for multi-line text within a card
def add_text_to_card(card_shape, text_lines, size=FONT_BODY_SIZE, color=TEXT_LIGHT, bold=False, top_margin=0.1, left_margin=0.1):
    tf = card_shape.text_frame
    tf.clear() # Clear existing text if any
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.TOP # Align text to top of the shape

    for i, line in enumerate(text_lines):
        p = tf.add_paragraph()
        p.text = line
        p.font.size = size
        p.font.color.rgb = color
        p.font.bold = bold
        p.space_after = Pt(4)

    tf.margin_left = Inches(left_margin)
    tf.margin_right = Inches(left_margin)
    tf.margin_top = Inches(top_margin)
    tf.margin_bottom = Inches(top_margin)

# Helper for key metric cards
def add_key_metric_card(slide, left, top, w, h, label, value, value_color=TEXT_LIGHT):
    card = add_card(slide, left, top, w, h)
    tf = card.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE # Center text vertically

    # Label
    p_label = tf.add_paragraph()
    p_label.text = label
    p_label.font.size = Pt(14)
    p_label.font.color.rgb = TEXT_GRAY
    p_label.alignment = MSO_ANCHOR.BOTTOM # Align label to bottom of text frame
    p_label.space_after = Pt(2)

    # Value
    p_value = tf.add_paragraph()
    p_value.text = value
    p_value.font.size = FONT_KEY_METRIC_SIZE
    p_value.font.color.rgb = value_color
    p_value.font.bold = True
    p_value.alignment = MSO_ANCHOR.TOP # Align value to top of text frame

    tf.margin_left = Inches(0.2)
    tf.margin_right = Inches(0.2)
    tf.margin_top = Inches(0.1)
    tf.margin_bottom = Inches(0.1)

    return card

# Helper for timeline cards
def add_timeline_card(slide, left, top, w, h, title, content_lines, status_marker=""):
    card = add_card(slide, left, top, w, h)
    tf = card.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.TOP

    # Title
    p_title = tf.add_paragraph()
    p_title.text = f"{status_marker} {title}" if status_marker else title
    p_title.font.size = FONT_CARD_TITLE_SIZE
    p_title.font.color.rgb = GOLD if "生效" in title else TEXT_LIGHT
    p_title.font.bold = True
    p_title.space_after = Pt(6)

    # Content
    for line in content_lines:
        p_content = tf.add_paragraph()
        p_content.text = line
        p_content.font.size = FONT_BODY_SIZE
        p_content.font.color.rgb = TEXT_GRAY
        p_content.space_after = Pt(4)

    tf.margin_left = Inches(0.2)
    tf.margin_right = Inches(0.2)
    tf.margin_top = Inches(0.2)
    tf.margin_bottom = Inches(0.2)

    return card

# Helper for status cards (Page 7)
def add_status_card(slide, left, top, w, h, title, status_text, status_color=TEXT_LIGHT):
    card = add_card(slide, left, top, w, h)
    tf = card.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE

    # Title
    p_title = tf.add_paragraph()
    p_title.text = title
    p_title.font.size = FONT_CARD_TITLE_SIZE
    p_title.font.color.rgb = TEXT_LIGHT
    p_title.font.bold = True
    p_title.alignment = MSO_ANCHOR.TOP

    # Status
    p_status = tf.add_paragraph()
    p_status.text = status_text
    p_status.font.size = FONT_SUBTITLE_SIZE
    p_status.font.color.rgb = status_color
    p_status.font.bold = True
    p_status.alignment = MSO_ANCHOR.BOTTOM

    tf.margin_left = Inches(0.2)
    tf.margin_right = Inches(0.2)
    tf.margin_top = Inches(0.2)
    tf.margin_bottom = Inches(0.2)


# --- Page 1: Title Slide ---
add_title_slide(prs,
                "00983D 質押套利方案：增加被動現金流可行性評估",
                "龍九控股 2026 年度策略規劃",
                "2026-07-28")

# --- Page 2: 核心策略 (Executive Summary) ---
slide2 = add_content_slide(prs,
                           "轉貸投資策略：創造穩定利差，優化現金流結構",
                           "透過低成本質押借款，投入高配息 ETF，實現每月淨現金流增長。")

# Card 1 (策略總覽)
card2_1 = add_card(slide2, 0.7, 2.5, 8.6, 1.2)
add_text_to_card(card2_1, ["質押借款 (2-3%) → 買入 00983D (7.1%) → 賺取利差 (~4%)"], size=Pt(16), bold=True, top_margin=0.3, left_margin=0.3)

# Card 2 (關鍵效益) - split into two cards for better spacing
card2_2 = add_card(slide2, 0.7, 4.0, 4.2, 2.5)
add_text_to_card(card2_2, ["關鍵效益:", "月配息滾入再投資", "加速資產增長"], size=Pt(14), color=TEXT_LIGHT, top_margin=0.2, left_margin=0.2)
card2_2.text_frame.paragraphs[0].font.bold = True # Make "關鍵效益" bold
card2_2.text_frame.paragraphs[1].font.color.rgb = GREEN # Highlight key benefit
card2_2.text_frame.paragraphs[2].font.color.rgb = GREEN # Highlight key benefit


# Card 3 (目標)
card2_3 = add_card(slide2, 5.1, 4.0, 4.2, 2.5)
add_text_to_card(card2_3, ["策略目標:", "提升被動收入", "清償高息負債"], size=Pt(14), color=TEXT_LIGHT, top_margin=0.2, left_margin=0.2)
card2_3.text_frame.paragraphs[0].font.bold = True # Make "策略目標" bold
card2_3.text_frame.paragraphs[1].font.color.rgb = BLUE # Highlight key benefit
card2_3.text_frame.paragraphs[2].font.color.rgb = BLUE # Highlight key benefit


# --- Page 3: 00983D 基本資料 (Key Investment Profile) ---
slide3 = add_content_slide(prs,
                           "00983D 富邦複合收益 ETF：穩健債券型標的",
                           "低波動特性與優異配息率，提供套利策略堅實基礎。")

# Layout: 2 rows of 3 cards each (total 6 cards), then one more.
card_width = 2.8
card_height = 1.8
start_left = 0.7
start_top = 2.5
col_spacing = 0.2
row_spacing = 0.3

# Row 1
add_key_metric_card(slide3, start_left, start_top, card_width, card_height, "名稱", "主動富邦複合收益")
add_key_metric_card(slide3, start_left + card_width + col_spacing, start_top, card_width, card_height, "類型", "主動操作債券型ETF")
add_key_metric_card(slide3, start_left + 2*(card_width + col_spacing), start_top, card_width, card_height, "掛牌", "2025年10月")

# Row 2
add_key_metric_card(slide3, start_left, start_top + card_height + row_spacing, card_width, card_height, "目前價格", "10.13")
add_key_metric_card(slide3, start_left + card_width + col_spacing, start_top + card_height + row_spacing, card_width, card_height, "年化配息率", "~7.1%", value_color=GREEN)
add_key_metric_card(slide3, start_left + 2*(card_width + col_spacing), start_top + card_height + row_spacing, card_width, card_height, "一年波動範圍", "±3.7% (極低)", value_color=GREEN)

# Row 3 (single card)
add_key_metric_card(slide3, start_left + card_width + col_spacing, start_top + 2*(card_height + row_spacing), card_width, card_height, "存續期間", "3.8年")


# --- Page 4: 資金計畫：階段性部署 (Phased Funding Plan) ---
slide4 = add_content_slide(prs,
                           "三階段資金部署：逐步擴大質押套利部位",
                           "結合自有資金與轉貸效益，穩健提升投資規模。")

timeline_start_left = 0.7
timeline_start_top = 2.5
timeline_card_width = 2.8
timeline_card_height = 3.5

# Phase 1
add_timeline_card(slide4, timeline_start_left, timeline_start_top, timeline_card_width, timeline_card_height,
                  "第一階段 (7-8月)",
                  ["自有資金 311萬 (戰備金)", "分批買入 00983D (目標 50張, 約50萬)", "月配息 ~6,000元"], status_marker="✨")

# Arrow (simple text representation for now)
arrow_left = timeline_start_left + timeline_card_width
arrow_top = timeline_start_top + timeline_card_height / 2 - 0.2
arrow_shape = slide4.shapes.add_textbox(Inches(arrow_left), Inches(arrow_top), Inches(col_spacing), Inches(0.5))
arrow_tf = arrow_shape.text_frame
p_arrow = arrow_tf.add_paragraph()
p_arrow.text = "►"
p_arrow.font.color.rgb = TEXT_LIGHT
p_arrow.font.size = Pt(24)
p_arrow.alignment = MSO_ANCHOR.MIDDLE

# Phase 2
add_timeline_card(slide4, timeline_start_left + timeline_card_width + col_spacing*2, timeline_start_top, timeline_card_width, timeline_card_height,
                  "第二階段 (9-10月)",
                  ["國泰轉貸完成 (2.185%)", "清償高息保單借貸", "質押增額再投入"], status_marker="🔄")

# Arrow
arrow_left2 = timeline_start_left + 2*(timeline_card_width + col_spacing*2) - col_spacing
arrow_shape2 = slide4.shapes.add_textbox(Inches(arrow_left2), Inches(arrow_top), Inches(col_spacing), Inches(0.5))
arrow_tf2 = arrow_shape2.text_frame
p_arrow2 = arrow_tf2.add_paragraph()
p_arrow2.text = "►"
p_arrow2.font.color.rgb = TEXT_LIGHT
p_arrow2.font.size = Pt(24)
p_arrow2.alignment = MSO_ANCHOR.MIDDLE

# Phase 3
add_timeline_card(slide4, timeline_start_left + 2*(timeline_card_width + col_spacing*2) + col_spacing, timeline_start_top, timeline_card_width, timeline_card_height,
                  "第三階段 (10月後)",
                  ["築巢優利貸生效", "擴大部位至 100-150張", "月配息目標 12,000-18,000元"], status_marker="🚀")


# --- Page 5: 現金流預測與效益分析 (Cash Flow Projection & Benefit Analysis) ---
slide5 = add_content_slide(prs,
                           "優化現金流結構：每月淨收益顯著增長",
                           "精準估算質押成本與配息收入，確保策略盈利能力。")

# Left Column: Investment Scale & Monthly Distribution
left_col_width = 4.0
right_col_width = 5.0
col_start_left = 0.7
col_start_top = 2.5

# Card 1 (投資規模)
card5_1 = add_card(slide5, col_start_left, col_start_top, left_col_width, 4.5)
add_text_to_card(card5_1,
                 ["投資規模與月配息預估:",
                  "• 50張 → 月配 ~6,000 元",
                  "• 100張 → 月配 ~12,000 元",
                  "• 150張 → 月配 ~18,000 元"],
                 size=Pt(15), color=TEXT_LIGHT, top_margin=0.2, left_margin=0.2)
card5_1.text_frame.paragraphs[0].font.bold = True
card5_1.text_frame.paragraphs[1].font.color.rgb = GREEN
card5_1.text_frame.paragraphs[2].font.color.rgb = GREEN
card5_1.text_frame.paragraphs[3].font.color.rgb = GREEN


# Right Column: Benefit Estimation
right_col_left = col_start_left + left_col_width + 0.3

card5_2 = add_card(slide5, right_col_left, col_start_top, right_col_width, 4.5)
tf5_2 = card5_2.text_frame
tf5_2.clear()
tf5_2.word_wrap = True
tf5_2.vertical_anchor = MSO_ANCHOR.TOP
tf5_2.margin_left = Inches(0.2)
tf5_2.margin_right = Inches(0.2)
tf5_2.margin_top = Inches(0.2)
tf5_2.margin_bottom = Inches(0.2)

p_header = tf5_2.add_paragraph()
p_header.text = "效益估算 (以 50萬借款為例):"
p_header.font.size = Pt(16)
p_header.font.color.rgb = TEXT_LIGHT
p_header.font.bold = True
p_header.space_after = Pt(10)

# Cost
p_cost_label = tf5_2.add_paragraph()
p_cost_label.text = "質押成本估算:"
p_cost_label.font.size = FONT_BODY_SIZE
p_cost_label.font.color.rgb = TEXT_GRAY
p_cost_label.space_after = Pt(2)

p_cost_value = tf5_2.add_paragraph()
p_cost_value.text = "年息 15,000 元 → 月息 1,250 元"
p_cost_value.font.size = FONT_CARD_TITLE_SIZE
p_cost_value.font.color.rgb = RED
p_cost_value.font.bold = True
p_cost_value.space_after = Pt(10)

# Net Cash Flow
p_net_label = tf5_2.add_paragraph()
p_net_label.text = "淨月現金流:"
p_net_label.font.size = FONT_BODY_SIZE
p_net_label.font.color.rgb = TEXT_GRAY
p_net_label.space_after = Pt(2)

p_net_value = tf5_2.add_paragraph()
p_net_value.text = "6,000 - 1,250 = 4,750 元"
p_net_value.font.size = FONT_CARD_TITLE_SIZE + Pt(4) # Larger for emphasis
p_net_value.font.color.rgb = GREEN
p_net_value.font.bold = True
p_net_value.space_after = Pt(10)

# Annual Net Gain
p_annual_label = tf5_2.add_paragraph()
p_annual_label.text = "年化淨收益估算:"
p_annual_label.font.size = FONT_BODY_SIZE
p_annual_label.font.color.rgb = TEXT_GRAY
p_annual_label.space_after = Pt(2)

p_annual_value = tf5_2.add_paragraph()
p_annual_value.text = "4,750 元/月 * 12 = 57,000 元/年"
p_annual_value.font.size = FONT_CARD_TITLE_SIZE
p_annual_value.font.color.rgb = GREEN
p_annual_value.font.bold = True


# --- Page 6: 風險控管與策略穩健性 (Risk Management & Strategy Robustness) ---
slide6 = add_content_slide(prs,
                           "多重風險控管機制：確保質押套利策略安全無虞",
                           "低波動標的與成本鎖定，有效規避市場與利率風險。")

card_start_left = 0.7
card_start_top = 2.5
card_w = 4.2
card_h = 2.0
col_gap = 0.3
row_gap = 0.4

# Card 1 (市場風險)
card6_1 = add_card(slide6, card_start_left, card_start_top, card_w, card_h)
add_text_to_card(card6_1, ["市場風險:", "00983D 波動極低 (±3.7%)", "→ 無斷頭風險"], size=Pt(14), color=TEXT_LIGHT, top_margin=0.2, left_margin=0.2)
card6_1.text_frame.paragraphs[0].font.bold = True
card6_1.text_frame.paragraphs[1].font.color.rgb = GREEN
card6_1.text_frame.paragraphs[2].font.color.rgb = GREEN

# Card 2 (利率風險)
card6_2 = add_card(slide6, card_start_left + card_w + col_gap, card_start_top, card_w, card_h)
add_text_to_card(card6_2, ["利率風險:", "質押利率固定 (2.185%)", "→ 利差穩定"], size=Pt(14), color=TEXT_LIGHT, top_margin=0.2, left_margin=0.2)
card6_2.text_frame.paragraphs[0].font.bold = True
card6_2.text_frame.paragraphs[1].font.color.rgb = GREEN
card6_2.text_frame.paragraphs[2].font.color.rgb = GREEN

# Card 3 (流動性風險)
card6_3 = add_card(slide6, card_start_left, card_start_top + card_h + row_gap, card_w, card_h)
add_text_to_card(card6_3, ["流動性風險:", "每月配息可覆蓋利息", "→ 財務彈性高"], size=Pt(14), color=TEXT_LIGHT, top_margin=0.2, left_margin=0.2)
card6_3.text_frame.paragraphs[0].font.bold = True
card6_3.text_frame.paragraphs[1].font.color.rgb = GREEN
card6_3.text_frame.paragraphs[2].font.color.rgb = GREEN

# Card 4 (操作風險)
card6_4 = add_card(slide6, card_start_left + card_w + col_gap, card_start_top + card_h + row_gap, card_w, card_h)
add_text_to_card(card6_4, ["操作風險:", "不做槓桿超額", "→ 維持保守策略"], size=Pt(14), color=TEXT_LIGHT, top_margin=0.2, left_margin=0.2)
card6_4.text_frame.paragraphs[0].font.bold = True
card6_4.text_frame.paragraphs[1].font.color.rgb = GREEN
card6_4.text_frame.paragraphs[2].font.color.rgb = GREEN


# --- Page 7: 進度狀態與下一步行動 (Current Status & Next Steps) ---
slide7 = add_content_slide(prs,
                           "策略執行進度：關鍵里程碑穩步推進",
                           "掌握各項子任務節點，確保計畫按時交付。")

status_card_start_left = 0.7
status_card_start_top = 2.5
status_card_w = 4.2
status_card_h = 2.0

# Status Card 1
add_status_card(slide7, status_card_start_left, status_card_start_top, status_card_w, status_card_h,
                "00983D 方案B", "⏳ 慢補中", GOLD)

# Status Card 2
add_status_card(slide7, status_card_start_left + status_card_w + col_gap, status_card_start_top, status_card_w, status_card_h,
                "國泰轉貸", "⏳ 9/25", GOLD)

# Status Card 3
add_status_card(slide7, status_card_start_left, status_card_start_top + status_card_h + row_gap, status_card_w, status_card_h,
                "築巢優利貸 2.185%", "⏳ 10/1生效", GOLD)

# Status Card 4
add_status_card(slide7, status_card_start_left + status_card_w + col_gap, status_card_start_top + status_card_h + row_gap, status_card_w, status_card_h,
                "信貸套利", "⏸️ 待第一階段完成", BLUE)


# Save the presentation
output_path = r"C:/Users/bot/Desktop/longjiu_system/轉貸投資可行性評估_專業版.pptx"
prs.save(output_path)
print(f"Presentation saved to: {output_path}")
