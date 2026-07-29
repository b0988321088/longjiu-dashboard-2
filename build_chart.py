"""
大轉向資產配置策略 — 圖表版
減少文字，增加圖表
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION, XL_LABEL_POSITION
import json, os

BASE = os.path.dirname(os.path.abspath(__file__))
SNAP = json.load(open(f'{BASE}/snapshot.json'))
INS = SNAP.get('allianz_combined',0) + SNAP.get('firstjin_fl65_current_value',0)
SEC = SNAP.get('securities_total_market_value',0)
FUND = SNAP.get('fund_market_value',0)
CASH = SNAP.get('real_liquid_assets',0)
TOTAL = INS + SEC + FUND + CASH
p = SNAP.get('penetration',{}).get('actual_pct',{})

prs = Presentation()
prs.slide_width = Inches(13.33); prs.slide_height = Inches(7.5)
BG = RGBColor(0x0B,0x0D,0x1A); WHITE = RGBColor(0xFF,0xFF,0xFF)
GRAY = RGBColor(0x8A,0x8F,0xA0); GOLD = RGBColor(0xF7,0xA0,0x1C)
GREEN = RGBColor(0x34,0xD3,0x99); RED = RGBColor(0xFF,0x5C,0x5C)

def ns():
    s = prs.slides.add_slide(prs.slide_layouts[6])
    s.background.fill.solid(); s.background.fill.fore_color.rgb = BG
    return s

def txt(s, t, l, tp, w, h, sz=14, bold=False, c=WHITE):
    tb = s.shapes.add_textbox(Inches(l), Inches(tp), Inches(w), Inches(h))
    p = tb.text_frame.paragraphs[0]
    p.text = t; p.font.size = Pt(sz); p.font.bold = bold; p.font.color.rgb = c

def chart(s, data, l, tp, w, h, typ=XL_CHART_TYPE.DOUGHNUT):
    c = s.shapes.add_chart(typ, Inches(l), Inches(tp), Inches(w), Inches(h), data).chart
    c.has_legend = True; c.legend.position = XL_LEGEND_POSITION.BOTTOM
    c.legend.include_in_layout = False
    c.legend.font.color.rgb = WHITE
    c.legend.font.name = 'Microsoft JhengHei'

    # Set font for data labels
    for series in c.series:
        series.has_data_labels = True # Enable data labels for all series
        data_labels = series.data_labels
        data_labels.font.name = 'Microsoft JhengHei'
        data_labels.font.color.rgb = WHITE
        # For doughnut charts, ensure data labels are shown as category name and percentage
        if typ == XL_CHART_TYPE.DOUGHNUT:
            data_labels.show_category_name = True
            data_labels.show_percentage = True
            data_labels.separator = '\n' # Add a newline separator if needed
            data_labels.position = XL_LABEL_POSITION.OUTSIDE_END

    # Set font for category axis labels (for bar/column charts)
    if typ in (XL_CHART_TYPE.BAR_CLUSTERED, XL_CHART_TYPE.COLUMN_CLUSTERED):
        if hasattr(c, 'category_axis'): # Defensive check, though it should exist for these types
            c.category_axis.tick_labels.font.name = 'Microsoft JhengHei'
            c.category_axis.tick_labels.font.color.rgb = WHITE

    # 圖表區背景透明
    plot = c.plots[0]
    try:
        plot.gap_width = 100
    except: pass
    return c

# === S1: 封面 ===
s = ns()
txt(s, '大轉向前的資產配置策略', 1.5, 1.5, 10, 1.2, 36, True, WHITE)
txt(s, 'USD/TWD 32.38 ｜ 台股崩跌 2,000 點 ｜ 10yr 4.60%', 1.5, 2.8, 10, 0.5, 18, False, GOLD)
txt(s, '龍九控股 · 2026-07-29 ｜ Chief Secretary + CIO', 1.5, 4.5, 8, 0.5, 14, False, GRAY)

# === S2: 配置圓餅圖 ===
s = ns()
txt(s, '目前資產配置', 0.5, 0.3, 5, 0.6, 28, True, WHITE)
txt(s, '台股不足20pp + 現金過多14pp = 雙重問題', 0.5, 1, 10, 0.4, 14, False, GRAY)
cd = CategoryChartData()
cd.categories = ['台股成長 15.2%','美股成長 33.7%','防守配息 12.4%','債券 19.4%','現金 19.2%']
cd.add_series('', (15.2, 33.7, 12.4, 19.4, 19.2))
c = chart(s, cd, 0.3, 1.5, 5.5, 5)
# 目標對照
txt(s, '目標配置對照', 6.5, 0.3, 5, 0.6, 28, True, WHITE)
cd2 = CategoryChartData()
cd2.categories = ['台股 35%','美股 30%','防守 25%','債券 5%','現金 5%']
cd2.add_series('', (35, 30, 25, 5, 5))
chart(s, cd2, 6.5, 1.5, 5.5, 5)
txt(s, '→ 台股需 +65 萬  ｜  現金須 -151 萬', 6.5, 6.5, 6, 0.4, 16, True, GOLD)

# === S3: ETF 比較長條圖 ===
s = ns()
txt(s, 'ETF 持股分析', 0.5, 0.3, 5, 0.6, 28, True, WHITE)
txt(s, '今日崩跌後：本益比 ~14 倍，極度便宜', 0.5, 1, 10, 0.4, 14, False, GRAY)
cd3 = CategoryChartData()
cd3.categories = ['0050','006208','00878','009816','00919','00713','0056']
cd3.add_series('持股市值(萬)', (20.1, 46.6, 48.8, 26.4, 17.8, 12.2, 5.0))
chart(s, cd3, 0.5, 1.5, 7, 4.5, XL_CHART_TYPE.COLUMN_CLUSTERED)
txt(s, '✅ 不動：0050/006208/00878/009816/00713/00919', 1, 6.2, 10, 0.3, 14, False, GREEN)
txt(s, '⚠️ 觀察：0056(凍結) ｜ ❌ 不買VT（安聯已全球分散）', 1, 6.6, 10, 0.3, 14, False, GOLD)
txt(s, '💡 加碼：00983D +70萬 ｜ 00878/00919 +30萬 ｜ 009816 +20萬', 1, 7, 10, 0.3, 14, True, WHITE)

# === S4: 債務結構 ===
s = ns()
txt(s, '債務結構 vs 市場利率 6-7%', 0.5, 0.3, 8, 0.6, 28, True, WHITE)
txt(s, '您的平均 ~2.5%，比市場低 3.5pp', 0.5, 1, 8, 0.4, 14, False, GRAY)
cd4 = CategoryChartData()
cd4.categories = ['保單借貸\n5%','市場房貸\n~6.5%','理財型\n4.0%','質押\n3.9%','築巢\n2.185%','國泰週轉\n2.6%']
cd4.add_series('利率(%)', (5.0, 6.5, 4.0, 3.9, 2.185, 2.6))
chart(s, cd4, 0.5, 1.5, 7, 4.5, XL_CHART_TYPE.COLUMN_CLUSTERED)
txt(s, '🏠 您的築巢 2.185% 比市場低 4.3pp', 1, 6.2, 10, 0.3, 16, True, GREEN)
txt(s, '🔴 4筆保單借貸@5% → 9月清償年省20萬', 1, 6.6, 10, 0.3, 16, False, RED)
txt(s, '💰 年省利息約85萬 vs 市場利率', 1, 7, 10, 0.3, 16, False, GOLD)

# === S5: 匯率 ===
s = ns()
txt(s, '匯率 USD/TWD：結構性弱勢保護', 0.5, 0.3, 8, 0.6, 28, True, WHITE)
txt(s, '81% 美元曝險 = 台幣貶值時您的資產自動增值', 0.5, 1, 10, 0.4, 14, False, GRAY)
cd5 = CategoryChartData()
cd5.categories = ['台幣28','台幣29','台幣30','台幣31','台幣32(目前)','台幣33']
vals = []
for r in [28,29,30,31,32,33]:
    vals.append(int(13160000/29.78 * r))  # USD exposure revalued
cd5.add_series('美元資產價值(萬)', [int(v/10000) for v in vals])
chart(s, cd5, 0.5, 1.5, 7, 4.5, XL_CHART_TYPE.COLUMN_CLUSTERED)
txt(s, '✅ 目前 32.38 = 結構性優勢', 1, 6.2, 10, 0.3, 16, True, GREEN)
txt(s, '⚠️ 不需因匯率調整投資，長期持有即可', 1, 6.6, 10, 0.3, 14, False, GOLD)

# === S6: 三情境比較 ===
s = ns()
txt(s, '三種情境 × 您的準備', 0.5, 0.3, 8, 0.6, 28, True, WHITE)
txt(s, 'A:升息20% ｜ B:高原50% ⭐ ｜ C:降息30%', 0.5, 1, 10, 0.4, 14, False, GRAY)
cd6 = CategoryChartData()
cd6.categories = ['情境A：升息','情境B：高原','情境C：降息']
cd6.add_series('適合度', (80, 100, 85))
chart(s, cd6, 0.5, 1.5, 5, 4, XL_CHART_TYPE.COLUMN_CLUSTERED)
txt(s, '🔴 A 升息：00983D抗跌 + 160萬防守 = 80分', 0.5, 5.8, 6, 0.3, 13, False, WHITE)
txt(s, '⭐ B 高原：00983D 7.1% + 台股加碼 = 100分', 0.5, 6.2, 6, 0.3, 13, False, GOLD)
txt(s, '🟢 C 降息：160萬子彈 + 美股資產升值 = 85分', 0.5, 6.6, 6, 0.3, 13, False, GREEN)
txt(s, '🏠 三情境下房貸 2.41% 都不受影響', 0.5, 7, 6, 0.3, 14, True, GREEN)

# === S7: 時間表 ===
s = ns()
txt(s, '行動時間表', 0.5, 0.3, 5, 0.6, 28, True, WHITE)
txt(s, '今日崩跌 = 更好的買點，冷靜執行', 0.5, 1, 10, 0.4, 14, False, GOLD)
cd7 = CategoryChartData()
cd7.categories = ['本週行動','7~8月','9月轉貸','Q4起']
cd7.add_series('投入金額(萬)', (50, 70, 65, 10))
chart(s, cd7, 0.5, 1.5, 7, 4, XL_CHART_TYPE.COLUMN_CLUSTERED)
txt(s, '📅 本週：00878+00919 +30萬 ｜ 009816/0050 +20萬', 1, 5.8, 10, 0.3, 14, False, WHITE)
txt(s, '📅 7-8月：00983D每週10張 +70萬（月配7K）', 1, 6.2, 10, 0.3, 14, False, WHITE)
txt(s, '📅 9月：築巢清償保單年省20萬 + 台股+65萬', 1, 6.6, 10, 0.3, 14, False, WHITE)
txt(s, '📅 Q4：現金160萬戰備 ｜ 年化被動收入+15~18萬', 1, 7, 10, 0.3, 14, False, GREEN)

# === S8: 結論 ===
s = ns()
txt(s, '結論：三層保護傘', 1, 0.5, 8, 0.8, 32, True, WHITE)
txt(s, '台幣弱勢保護 + 低利房貸紅利 + 充足現金子彈', 1, 1.3, 10, 0.5, 18, False, GOLD)
txt(s, '🔥 大轉向來臨時：升息有匯率保護 ｜ 高原有配息收 ｜ 降息有子彈加', 1, 2.5, 11, 0.5, 18, True, WHITE)
txt(s, '年化增加被動收入 15~18 萬  ｜  22 個月安全期  ｜  不用恐慌', 1, 3.5, 11, 0.5, 16, False, GREEN)

prs.save(f'{BASE}/大轉向圖表版.pptx')
print('✅ 圖表版完成！8頁（含5張圖表）')
