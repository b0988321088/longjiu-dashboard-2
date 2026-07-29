from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

prs = Presentation()
prs.slide_width = Inches(10)
prs.slide_height = Inches(7.5)

BG = RGBColor(0x0B,0x0D,0x1A)
CARD = RGBColor(0x15,0x18,0x2E)
GOLD = RGBColor(0xF7,0xA0,0x1C)
GRN = RGBColor(0x34,0xD3,0x99)
RED = RGBColor(0xFF,0x5C,0x5C)
BLU = RGBColor(0x5B,0x9B,0xF7)
TXT = RGBColor(0xE8,0xE8,0xF0)
GRY = RGBColor(0x8A,0x8F,0xA0)
WHT = RGBColor(0xFF,0xFF,0xFF)

def set_bg(sl):
    sl.background.fill.solid()
    sl.background.fill.fore_color.rgb = BG

def card(sl,l,t,w,h):
    s=sl.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(l),Inches(t),Inches(w),Inches(h))
    s.fill.solid();s.fill.fore_color.rgb=CARD
    s.line.color.rgb=RGBColor(0x1E,0x22,0x3D);s.line.width=Pt(1)

def add_txt(sl,text,l,t,w,h,sz=14,clr=TXT,bold=False,align=PP_ALIGN.LEFT):
    tb=sl.shapes.add_textbox(Inches(l),Inches(t),Inches(w),Inches(h))
    tb.text_frame.word_wrap=True
    p=tb.text_frame.paragraphs[0];p.text=text
    p.font.size=Pt(sz);p.font.color.rgb=clr;p.font.bold=bold;p.alignment=align

def metric(sl,l,t,w,h,label,val,sub="",clr=GRN):
    card(sl,l,t,w,h)
    add_txt(sl,label,l+0.2,t+0.15,w-0.4,0.3,11,GRY)
    add_txt(sl,val,l+0.2,t+0.5,w-0.4,0.5,26,clr,True)
    if sub: add_txt(sl,sub,l+0.2,t+1.15,w-0.4,0.3,10,GRY)

def info_card(sl,l,t,w,h,lines,sz=12,clr=TXT,title="",tc=WHT):
    card(sl,l,t,w,h)
    y=t+0.2
    if title:
        add_txt(sl,title,l+0.3,y,w-0.6,0.4,14,tc,True)
        y+=0.5
    for ln in lines:
        if ln.startswith("**") and ln.endswith("**"):
            add_txt(sl,ln.strip("**"),l+0.3,y,w-0.6,0.3,13,clr,True)
            y+=0.35
        else:
            add_txt(sl,ln,l+0.3,y,w-0.6,0.3,sz,clr)
            y+=0.28

# === S1: Cover ===
s=prs.slides.add_slide(prs.slide_layouts[6]);set_bg(s)
card(s,0.8,1.5,8.4,4.5)
add_txt(s,"龍九控股",1.3,1.8,5,0.5,14,BLU)
add_txt(s,"轉貸投資可行性評估計劃",1.3,2.3,8,1,36,WHT,True)
add_txt(s,"以 2.185% 低成本資金活化資產結構，建立穩健套利現金流",1.3,3.3,8,0.5,14,GRY)
add_txt(s,"制定日期：2026-07-26 | 評估期間：2026 Q3~Q4",1.3,4.2,5,0.4,12,GRY)

# === S2: 資產負債 ===
s=prs.slides.add_slide(prs.slide_layouts[6]);set_bg(s)
add_txt(s,"現有資產負債總覽",0.5,0.2,9,0.5,24,WHT,True)
metric(s,0.5,0.9,2.2,1.5,"不動產","3,332 萬","佔比 65.7%",BLU)
metric(s,2.9,0.9,2.2,1.5,"保單現值","980 萬","佔比 19.3%",GRN)
metric(s,5.3,0.9,2.2,1.5,"證券+基金","327 萬","佔比 6.5%",GOLD)
metric(s,7.7,0.9,2.2,1.5,"現金","331 萬","佔比 6.5%",GRN)
info_card(s,0.5,2.7,4.5,2.2,["永豐房貸 x3   1,316 萬  ~2.5%","理財型房貸     301 萬  ~3.5%","保單借貸-安聯A  200 萬  ~5% 🔴","保單借貸-安聯B  100 萬  ~5% 🔴","保單借貸-第一金 100 萬  ~5% 🔴","證券質押       100 萬  ~3%","總負債 2,117 萬 | 負債比 41.8%"],sz=12,title="負債結構",tc=RED)
info_card(s,5.3,2.7,4.2,2.2,["月收入：206,859","  保單配息 73K / 房租 56K","  薪資 43K / 股息 10K","月支出：210,958","  房貸 99K / 信用卡 38K","  生活 45K / 利息 ~18.5K","⚠️ 月缺口：-4,099"],sz=12,title="每月收支",tc=GOLD)

# === S3: 方案比較 ===
s=prs.slides.add_slide(prs.slide_layouts[6]);set_bg(s)
add_txt(s,"轉貸方案比較",0.5,0.2,9,0.5,24,WHT,True)
info_card(s,0.5,1,2.9,3.5,["方案 A ⭐⭐⭐⭐⭐","國泰→築巢優利貸","利率：2.185%","月付：~49,800","總利息：~470 萬","月省房貸：~49,658"],sz=12,title="",tc=GRN)
info_card(s,3.6,1,2.9,3.5,["方案 B ⭐⭐","國泰轉貸維持","利率：2.6%","月付：~52,500","總利息：~560 萬"],sz=12,title="",tc=GRY)
info_card(s,6.7,1,2.9,3.5,["方案 C ⭐⭐","維持永豐現狀","利率：~2.5%","月付：99,458","總利息：~1,050 萬"],sz=12,title="",tc=GRY)
info_card(s,0.5,4.8,9,1.5,["結論：方案 A 最優，月省 ~49,658 房貸支出","轉貸後月付從 99,458 → 49,800，降幅 50%","年省利息 ~192,000（保單質押 5%→0%）"],sz=13,title="",tc=GOLD)

# === S4: 資金流向 ===
s=prs.slides.add_slide(prs.slide_layouts[6]);set_bg(s)
add_txt(s,"資金來源與流向",0.5,0.2,9,0.5,24,WHT,True)
info_card(s,0.5,1,9,2.5,["國泰轉貸撥款（9 月底）","├── 清償永豐房貸 1,316 萬（9/25 到期）","├── 清償保單借貸 400 萬（利率 ~5%→0%）","│     ├─ 安聯A 200 萬 │ ├─ 安聯B 100 萬","│     └─ 第一金FL65 100 萬","└── 剩餘資金 ~300 萬（低利 2.185% 可用於投資）"],sz=12,title="",tc=TXT)
info_card(s,0.5,3.8,9,2.5,["轉貸後每月現金流變化","房貸：99,458 → 49,800（-49,658）✅","保單質押利息：~16,000 → 0 ✅","保單配息：73,000 → 89,000（+16,000）✅","月缺口：-4,099 → +61,559 ✅"],sz=12,title="",tc=GRN)

# === S5: 風險 ===
s=prs.slides.add_slide(prs.slide_layouts[6]);set_bg(s)
add_txt(s,"風險評估",0.5,0.2,9,0.5,24,WHT,True)
risks=[("利率上升","固定利率 2.185%\n已鎖定，不受影響",GRN),("核貸失敗","國泰已面簽/對保 ✅",GRN),
       ("房價下跌","負債比 41.8%\n安全邊際充足",GOLD),("投資虧損","分批買入\n配息保護",GOLD),
       ("解約損失","無提前還款\n違約金 ✅",GRN)]
for i,(nm,desc,clr) in enumerate(risks):
    l=0.5+i*1.9
    info_card(s,l,1,1.7,2.5,[f"🔴 {nm}","",desc],sz=11,title="",tc=clr)
info_card(s,0.5,3.8,9,2.5,["第一階段（撥款後立即）：清償高利負債 400 萬→年省利息~19.2 萬","第二階段（撥款後 1 週內）：建立 300 萬安全網→戰備金+低點布局","第三階段（10 月後）：擴大 00983D 質押套利→增加月配~6,000"],sz=12,title="三階段資金部署",tc=GOLD)

# === S6: 三階段 ===
s=prs.slides.add_slide(prs.slide_layouts[6]);set_bg(s)
add_txt(s,"三階段資金部署計劃",0.5,0.2,9,0.5,24,WHT,True)
stages=[("第一階段：清償高利","撥款後立即執行",["✅ 清償安聯A 200 萬","✅ 清償安聯B 100 萬","✅ 清償第一金 100 萬","💰 年省利息 ~19.2 萬\n  配息全額領取"],GRN),
        ("第二階段：安全網","撥款後 1 週內",["💰 300 萬戰備金","🏠 預留 150 萬房貸","📈 100 萬 00983D 分批","💵 50 萬靈活運用"],BLU),
        ("第三階段：擴大套利","10 月後",["🎯 擴大 00983D 至 150 張","💵 月配 12K~18K","📈 盈餘 8.2 萬續投","🇹🇼 台股逢低補碼"],GOLD)]
for i,(tt,tm,items,clr) in enumerate(stages):
    l=0.5+i*3.2
    info_card(s,l,1,2.9,3.8,items,sz=11,title=tt,tc=clr)
    add_txt(s,tm,l+0.2,1.4,2.5,0.3,10,GRY)
info_card(s,0.5,5.1,9,1.5,["**每月現金流：-4,099 → +61,559**","**每年增被動收入：12.8~25.4 萬 | 負債比 41.8%→~30%**"],sz=13,clr=GRN,title="最終成效",tc=GRN)

# === S7: End ===
s=prs.slides.add_slide(prs.slide_layouts[6]);set_bg(s)
card(s,1.5,2,7,3.5)
add_txt(s,"龍九控股",2.5,2.5,5,0.5,14,BLU,align=PP_ALIGN.CENTER)
add_txt(s,"轉貸投資可行性評估計劃",2,3,6,1,32,WHT,True,PP_ALIGN.CENTER)
add_txt(s,"制定日期：2026-07-26 | Chief Secretary",2.5,4.2,5,0.4,12,GRY,align=PP_ALIGN.CENTER)

prs.save('C:/Users/bot/Desktop/longjiu_system/轉貸投資計劃.pptx')
print('✅ OK')
